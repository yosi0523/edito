# -*- coding: utf-8 -*-
"""
Tenneco × Hyundai Project — Toss-style Korean Presentation (18 slides)
For Tenneco US HQ internal seminar — 16:9

말풍선(의뢰인) 요청 전면 반영:
- 세계지도 + 국가별 노란 원 하이라이트
- 국기 + 차량사진 자리 + 프로젝트명 카드
- "Sales by country infographics" 스타일 랭킹
- 북미·남미 지도에 주(state)/도시 이름만 심플 표시
- USA·Brazil 강조 하이라이트 지도
- 북미 ECS 공급 계획 심플 표
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# ===== Tenneco brand color system (Toss-style layout + Tenneco brand colors) =====
# Reference: Tenneco logo royal blue + deep navy background
# Variable names kept for compatibility; values updated to Tenneco brand palette.
TEAL       = RGBColor(0x1E, 0x40, 0xAF)   # Tenneco Royal Blue  (primary)
TEAL_DEEP  = RGBColor(0x0F, 0x2E, 0x61)   # Tenneco Navy        (deep emphasis)
TEAL_SOFT  = RGBColor(0xEF, 0xF6, 0xFF)   # Very light blue     (card background)
TEAL_MID   = RGBColor(0x3B, 0x82, 0xF6)   # Mid blue            (secondary)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)  # Bright blue         (highlight)
LIME       = RGBColor(0xFA, 0xCC, 0x15)   # Gold yellow         (call-outs)
YELLOW_HL  = RGBColor(0xFD, 0xE0, 0x47)   # Brighter yellow     (country highlights on maps)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN      = RGBColor(0x05, 0x96, 0x69)

TEXT_BLACK   = RGBColor(0x19, 0x1F, 0x28)
TEXT_GRAY    = RGBColor(0x6B, 0x76, 0x84)
TEXT_LIGHT   = RGBColor(0x8B, 0x95, 0xA1)
BORDER_GRAY  = RGBColor(0xE5, 0xE8, 0xEB)
BG_TINT      = RGBColor(0xF9, 0xFA, 0xFB)
BG_GRAY      = RGBColor(0xF2, 0xF4, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DOT_BG       = RGBColor(0xCF, 0xD8, 0xDC)   # world map base dots

FONT_KR = "맑은 고딕"
FONT_EN = "Arial"

# ===== Presentation setup =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# =======================================================================
#                             H E L P E R S
# =======================================================================
def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bgr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bgr.shadow.inherit = False
    bgr.fill.solid()
    bgr.fill.fore_color.rgb = bg
    bgr.line.fill.background()
    return s


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = shadow
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w) if line_w else Pt(0.75)
    return shp


def rounded(slide, x, y, w, h, fill=None, line=None, radius=0.08,
            line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def oval(slide, x, y, w, h, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w) if line_w else Pt(0.75)
    return shp


def txt(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_BLACK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR,
        italic=False, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def line(slide, x1, y1, x2, y2, color=BORDER_GRAY, weight=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def hairline(slide, x, y, w, color=BORDER_GRAY, weight=0.75):
    return line(slide, x, y, x + w, y, color=color, weight=weight)


def eyebrow(slide, x, y, text, color=TEXT_GRAY, size=10):
    """Small uppercase eyebrow label — Toss signature."""
    txt(slide, x, y, Inches(5), Inches(0.25),
        text, size=size, bold=True, color=color,
        font=FONT_EN)


def page_num(slide, n, total=18):
    txt(slide, Inches(12.5), Inches(7.15), Inches(0.55), Inches(0.25),
        f"{n:02d} / {total:02d}", size=9, color=TEXT_LIGHT,
        align=PP_ALIGN.RIGHT, font=FONT_EN)


def footer(slide, n=None):
    txt(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25),
        "General Business – Tenneco Confidential",
        size=9, color=TEXT_LIGHT, italic=True)
    if n is not None:
        page_num(slide, n)


def slide_title(slide, title, sub=None, y_title=0.55):
    """Standard Toss-style title block."""
    txt(slide, Inches(0.6), Inches(y_title), Inches(12), Inches(0.65),
        title, size=28, bold=True, color=TEXT_BLACK)
    if sub:
        txt(slide, Inches(0.6), Inches(y_title + 0.75),
            Inches(12), Inches(0.32),
            sub, size=13, color=TEXT_GRAY)


# ---------- World map (dot-based stylized world map) ----------
def world_map_dots(slide, x, y, w, h, dot_color=DOT_BG, bg=None):
    """Draw a dot-pattern world map (stylized, similar to reference images).
       Uses small circles arranged to roughly outline continents."""
    if bg is not None:
        rect(slide, x, y, w, h, fill=bg)
    # Each point is (rel_x, rel_y) normalized 0~1 inside map area
    # Rough continent silhouettes via scattered dots
    pts = []

    def grid(xmin, xmax, ymin, ymax, cols, rows):
        for i in range(cols):
            for j in range(rows):
                fx = xmin + (xmax - xmin) * (i / max(cols - 1, 1))
                fy = ymin + (ymax - ymin) * (j / max(rows - 1, 1))
                pts.append((fx, fy))

    # North America
    grid(0.07, 0.22, 0.22, 0.40, 8, 5)
    grid(0.09, 0.20, 0.40, 0.48, 4, 2)
    # Central / Mexico
    grid(0.13, 0.19, 0.48, 0.55, 3, 2)
    # South America
    grid(0.18, 0.26, 0.58, 0.80, 4, 6)
    grid(0.20, 0.24, 0.80, 0.90, 2, 3)
    # Europe
    grid(0.44, 0.56, 0.25, 0.38, 7, 4)
    # Africa
    grid(0.48, 0.58, 0.42, 0.68, 5, 6)
    # Middle East
    grid(0.55, 0.62, 0.38, 0.48, 4, 3)
    # Asia
    grid(0.60, 0.80, 0.22, 0.40, 10, 5)
    # India
    grid(0.65, 0.71, 0.42, 0.52, 3, 3)
    # SE Asia
    grid(0.73, 0.82, 0.44, 0.54, 4, 3)
    # Australia
    grid(0.80, 0.88, 0.70, 0.80, 4, 3)

    dot_r = Emu(int(min(w, h) * 0.008))
    for rx, ry in pts:
        cx = x + Emu(int(w * rx))
        cy = y + Emu(int(h * ry))
        oval(slide, cx - dot_r, cy - dot_r, dot_r * 2, dot_r * 2,
             fill=dot_color)


def highlight_country(slide, map_x, map_y, map_w, map_h, rx, ry,
                      size=0.55, color=YELLOW_HL, label=None,
                      label_color=TEXT_BLACK, label_size=10):
    """Draw a highlight circle over a country position (normalized coords)."""
    cx = map_x + Emu(int(map_w * rx))
    cy = map_y + Emu(int(map_h * ry))
    r = Inches(size / 2)
    # soft outer glow
    outer = oval(slide, cx - r - Inches(0.08), cy - r - Inches(0.08),
                 Inches(size + 0.16), Inches(size + 0.16),
                 fill=color)
    # Try to set transparency via XML for outer glow
    try:
        sp = outer.fill.fore_color._xFill
        alpha = etree.SubElement(sp, qn('a:alpha'))
        alpha.set('val', '30000')
    except Exception:
        pass
    # main circle
    oval(slide, cx - r, cy - r, Inches(size), Inches(size),
         fill=color, line=WHITE, line_w=1.2)
    if label:
        txt(slide, cx - Inches(0.9), cy - Inches(0.12),
            Inches(1.8), Inches(0.25),
            label, size=label_size, bold=True, color=label_color,
            align=PP_ALIGN.CENTER, font=FONT_EN)


def numbered_marker(slide, x, y, num, color=ACCENT_BLUE,
                    size=0.35, text_color=WHITE):
    """Numbered pin marker (for ranked infographic)."""
    r = Inches(size / 2)
    oval(slide, x - r, y - r, Inches(size), Inches(size),
         fill=color, line=WHITE, line_w=1.5)
    txt(slide, x - r, y - r, Inches(size), Inches(size),
        str(num), size=11, bold=True, color=text_color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)


# ---------- Flag helper ----------
FLAGS = {
    "CN": [RGBColor(0xDE, 0x28, 0x10)],
    "IN": [RGBColor(0xFF, 0x99, 0x33),
           RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x13, 0x88, 0x08)],
    "BR": [RGBColor(0x00, 0x9C, 0x3B),
           RGBColor(0xFE, 0xDF, 0x00)],
    "KR": [RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x00, 0x20, 0x71)],
    "US": [RGBColor(0xB2, 0x22, 0x34),
           RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x3C, 0x3B, 0x6E)],
    "DE": [RGBColor(0x00, 0x00, 0x00),
           RGBColor(0xDD, 0x00, 0x00),
           RGBColor(0xFF, 0xCE, 0x00)],
    "MX": [RGBColor(0x00, 0x69, 0x33),
           RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0xCE, 0x11, 0x26)],
    "CA": [RGBColor(0xFF, 0x00, 0x00),
           RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0xFF, 0x00, 0x00)],
    "AR": [RGBColor(0x74, 0xAC, 0xDF),
           RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x74, 0xAC, 0xDF)],
}


def flag(slide, x, y, w, h, code):
    colors = FLAGS.get(code, [TEXT_LIGHT])
    # outer
    border = rect(slide, x, y, w, h, fill=WHITE,
                  line=BORDER_GRAY, line_w=0.5)
    n = len(colors)
    stripe_h = Emu(int(h / n))
    for i, c in enumerate(colors):
        rect(slide, x, y + Emu(int(h * i / n)), w, stripe_h, fill=c)


# ---------- Vehicle card (flag + car placeholder + project code + name) ----------
def vehicle_card(slide, x, y, w, h, code, country, project_code, vehicle,
                 year=None, new_badge=False, highlight=False):
    """Card with country flag, car placeholder, project code, vehicle name.
       Toss-style: white bg, soft border, rounded."""
    fill_c = TEAL_SOFT if highlight else WHITE
    rounded(slide, x, y, w, h,
            fill=fill_c, line=BORDER_GRAY, line_w=0.75, radius=0.08)
    # Flag + country (top row)
    flag(slide, x + Inches(0.15), y + Inches(0.15),
         Inches(0.32), Inches(0.2), code)
    txt(slide, x + Inches(0.52), y + Inches(0.12),
        Inches(1.5), Inches(0.25),
        country, size=9, bold=True, color=TEXT_GRAY, font=FONT_EN)
    if new_badge:
        rounded(slide, x + w - Inches(0.65), y + Inches(0.13),
                Inches(0.5), Inches(0.22), fill=TEAL, radius=0.5)
        txt(slide, x + w - Inches(0.65), y + Inches(0.13),
            Inches(0.5), Inches(0.22),
            "NEW", size=7, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # Car placeholder (middle)
    photo_y = y + Inches(0.4)
    photo_h = h - Inches(1.05)
    rect(slide, x + Inches(0.15), photo_y, w - Inches(0.3),
         photo_h, fill=BG_GRAY)
    txt(slide, x + Inches(0.15), photo_y, w - Inches(0.3), photo_h,
        "🚗 차량사진", size=9, color=TEXT_LIGHT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=True)
    # Project code + vehicle
    label_y = y + h - Inches(0.55)
    txt(slide, x + Inches(0.15), label_y, w - Inches(0.3), Inches(0.28),
        project_code, size=13, bold=True, color=TEXT_BLACK, font=FONT_EN)
    line2 = vehicle + (f" · {year}" if year else "")
    txt(slide, x + Inches(0.15), label_y + Inches(0.28),
        w - Inches(0.3), Inches(0.23),
        line2, size=10, color=TEXT_GRAY)


# ---------- Simple pie chart using shapes (for GM/Toyota/Hyundai etc.) ----------
def status_pill(slide, x, y, w, h, text, tone="red"):
    """Toss-style pill tag for status (Sales Decline / Market Share ↑)."""
    color = {"red": RED, "green": GREEN,
             "blue": ACCENT_BLUE, "gray": TEXT_GRAY}[tone]
    rounded(slide, x, y, w, h, fill=color, radius=0.5)
    txt(slide, x, y, w, h, text, size=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def stat_bar(slide, x, y, total_w, h, pct, fill_color=TEAL,
             bg_color=BG_GRAY):
    """Horizontal progress bar (Toss-style)."""
    rounded(slide, x, y, total_w, h, fill=bg_color, radius=0.5)
    filled_w = Emu(int(total_w * pct / 100))
    if filled_w > 0:
        rounded(slide, x, y, filled_w, h, fill=fill_color, radius=0.5)


def stat_stackbar(slide, x, y, w, h, segments):
    """Multi-segment stacked horizontal bar.
       segments: list of (pct, color, label). pct values sum to 100."""
    # Background
    rounded(slide, x, y, w, h, fill=BG_GRAY, radius=0.3)
    cursor = x
    for pct, color, _ in segments:
        seg_w = Emu(int(w * pct / 100))
        if seg_w <= 0:
            continue
        rect(slide, cursor, y, seg_w, h, fill=color)
        cursor = cursor + seg_w


def add_donut_chart(slide, x, y, w, h, data, title=None):
    """Real donut chart using python-pptx chart API."""
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in data]
    chart_data.add_series("Share", [d[1] for d in data])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, w, h, chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    # data label
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = True
    dl.show_category_name = False
    dl.show_value = False
    dl.font.size = Pt(11)
    dl.font.bold = True
    # series colors
    series = chart.series[0]
    series_colors = [TEAL, TEAL_MID, TEXT_LIGHT, RED]
    for i, _ in enumerate(data):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = series_colors[i % len(series_colors)]
    return chart


# ---------- Clean table (Toss-style: no vertical borders, hairline rows) ----------
def toss_table(slide, x, y, w, h, data, col_widths=None,
               header_fill=None, header_color=TEXT_GRAY,
               header_size=11, body_size=12,
               highlight_row=None, highlight_fill=TEAL_SOFT,
               body_color=TEXT_BLACK,
               first_col_bold=True, center_cols=None,
               row_heights=None):
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * cw / total))
    if row_heights:
        for i, rh in enumerate(row_heights):
            tbl.rows[i].height = Inches(rh)
    center_cols = center_cols or []
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            # Background
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill if header_fill else WHITE
            elif highlight_row is not None and ri == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            # Text
            tf = cell.text_frame
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = ""
            if ci in center_cols or ri == 0:
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT_KR
            if ri == 0:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = header_color
            else:
                r.font.size = Pt(body_size)
                if ci == 0 and first_col_bold:
                    r.font.bold = True
                    r.font.color.rgb = TEXT_BLACK
                else:
                    r.font.color.rgb = body_color
    return tbl


# =======================================================================
#                              S L I D E S
# =======================================================================

# -------- Slide 1 — Cover --------
def slide_01_cover():
    s = add_slide(WHITE)
    # tiny eyebrow top-left
    eyebrow(s, Inches(0.6), Inches(0.55),
            "TENNECO × HYUNDAI  ·  2026", color=TEAL, size=11)
    # Big title block left
    txt(s, Inches(0.6), Inches(2.6), Inches(10), Inches(0.7),
        "Hyundai Motor Group",
        size=38, bold=False, color=TEXT_GRAY)
    txt(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.4),
        "프로젝트 현황\n공략 전략",
        size=58, bold=True, color=TEXT_BLACK, spacing=1.1)
    # Thin teal bar + subtitle below
    rect(s, Inches(0.6), Inches(5.7), Inches(0.5), Inches(0.04),
         fill=TEAL)
    txt(s, Inches(0.6), Inches(5.8), Inches(10), Inches(0.35),
        "Internal Seminar  ·  US Headquarters",
        size=14, color=TEXT_GRAY, font=FONT_EN)
    # Bottom right
    txt(s, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
        "2026.04", size=10, color=TEXT_LIGHT,
        align=PP_ALIGN.RIGHT, font=FONT_EN)
    # Subtle decoration - a teal dot pattern accent
    for i in range(8):
        for j in range(4):
            x_ = Inches(9.5 + i * 0.35)
            y_ = Inches(0.8 + j * 0.35)
            r_ = Inches(0.05)
            oval(s, x_, y_, r_ * 2, r_ * 2, fill=TEAL_SOFT)


# -------- Slide 2 — 목차 --------
def slide_02_toc():
    s = add_slide(WHITE)
    eyebrow(s, Inches(0.6), Inches(0.6), "AGENDA", color=TEAL)
    txt(s, Inches(0.6), Inches(0.9), Inches(6), Inches(0.8),
        "목차", size=44, bold=True, color=TEXT_BLACK)
    # Divider
    hairline(s, Inches(0.6), Inches(2.0), Inches(8), color=BORDER_GRAY)

    items = [
        ("01", "현대 글로벌 시장 판매",       "Hyundai Global Market Share"),
        ("02", "현대 북미 및 남미",           "N.America & S.America Market"),
        ("03", "현대에서의 테네코 점유율",     "Tenneco Business Ratio in HMC"),
        ("04", "테네코에서의 현대 공급 프로그램", "Supply Programs  ·  CN · IN · BR"),
        ("05", "경쟁사 비교 내용",            "Competitor Analysis"),
        ("06", "현대 공략 방법",              "Attack Strategy"),
        ("07", "Q&A",                        "Discussion"),
    ]
    y0 = Inches(2.3)
    row_h = Inches(0.6)
    for i, (num, kr, en) in enumerate(items):
        y = y0 + row_h * i
        txt(s, Inches(0.6), y, Inches(0.9), Inches(0.5),
            num, size=22, bold=True, color=TEAL, font=FONT_EN)
        txt(s, Inches(1.6), y + Inches(0.03), Inches(7), Inches(0.4),
            kr, size=18, bold=True, color=TEXT_BLACK)
        txt(s, Inches(1.6), y + Inches(0.35), Inches(7), Inches(0.3),
            en, size=10, color=TEXT_GRAY, font=FONT_EN, italic=True)
    footer(s, 2)


# -------- Slide 3 — Hyundai Global Market Share 2025 --------
def slide_03_global_market():
    s = add_slide(WHITE)
    slide_title(s, "현대는 2025년, 414만 대를 팔았다",
                "Hyundai Global Market Share · 2025")

    # Hero stat — LEFT
    eyebrow(s, Inches(0.6), Inches(1.9),
            "2025 GLOBAL SALES", color=TEAL)
    txt(s, Inches(0.6), Inches(2.15), Inches(6), Inches(1.4),
        "414만 대", size=96, bold=True, color=TEAL, font=FONT_EN)
    txt(s, Inches(0.6), Inches(3.55), Inches(6), Inches(0.35),
        "Total Wholesale Volume",
        size=12, color=TEXT_GRAY, font=FONT_EN)
    # Commentary card
    rounded(s, Inches(0.6), Inches(4.15), Inches(5.8), Inches(1.4),
            fill=BG_TINT, line=None, radius=0.08)
    rect(s, Inches(0.6), Inches(4.15), Inches(0.08), Inches(1.4),
         fill=TEAL)
    txt(s, Inches(0.85), Inches(4.25), Inches(5.5), Inches(0.3),
        "Key insight", size=10, bold=True, color=TEAL, font=FONT_EN)
    txt(s, Inches(0.85), Inches(4.55), Inches(5.5), Inches(0.95),
        "지정학·경기 변동 속에서도\n"
        "북미 역대 최대 + 인도 견조한 성장으로\n"
        "글로벌 판매 비중 성공적 다변화",
        size=11, color=TEXT_BLACK, spacing=1.3)

    # Regional bars — RIGHT (말풍선: "이런 식의 그림..." → simplified bars)
    eyebrow(s, Inches(7.0), Inches(1.9),
            "BY REGION", color=TEAL)
    regions = [
        ("아시아",   "146만 대", 35.3, 100),
        ("아메리카", "134만 대", 32.4, 92),
        ("기타",     "73만 대",  17.6, 50),
        ("유럽",     "61만 대",  14.7, 42),
    ]
    y_bar = Inches(2.25)
    for i, (reg, val, pct, w_pct) in enumerate(regions):
        yy = y_bar + Inches(0.72 * i)
        # label
        txt(s, Inches(7.0), yy, Inches(1.3), Inches(0.28),
            reg, size=13, bold=True, color=TEXT_BLACK)
        # value + pct
        txt(s, Inches(10.5), yy, Inches(2.3), Inches(0.28),
            f"{val}  ·  {pct}%", size=11, color=TEXT_GRAY,
            align=PP_ALIGN.RIGHT, font=FONT_EN)
        # bar
        stat_bar(s, Inches(7.0), yy + Inches(0.35), Inches(5.8),
                 Inches(0.2), w_pct, fill_color=TEAL,
                 bg_color=BG_GRAY)

    # Bottom — Tenneco Global scale (말풍선 "이런 식의 글도 들어가면")
    rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.02),
         fill=BORDER_GRAY)
    eyebrow(s, Inches(0.6), Inches(6.0),
            "TENNECO GLOBAL SCALE", color=TEXT_GRAY, size=10)

    stats = [
        ("60,000",  "Global Team Members"),
        ("184",     "Manufacturing Plants"),
        ("23",      "Distribution Centers"),
        ("40",      "Engineering Facilities"),
    ]
    col_w = Inches(3.0)
    for i, (num, lbl) in enumerate(stats):
        x = Inches(0.6) + col_w * i
        txt(s, x, Inches(6.3), col_w, Inches(0.55),
            num, size=32, bold=True, color=TEAL, font=FONT_EN)
        txt(s, x, Inches(6.88), col_w, Inches(0.28),
            lbl, size=11, color=TEXT_GRAY, font=FONT_EN)
    footer(s, 3)


# -------- Slide 4 — N.America & S.America --------
# 말풍선: "이런 식의 그림이 사용되어 표시되었으면" → 세계지도 + 국가 원형 마커
def slide_04_america():
    s = add_slide(WHITE)
    slide_title(s, "북미·남미 합계 134만 대",
                "Hyundai N.America & S.America Market · 2025")

    # LEFT — World map with country highlights (말풍선 요청 구현)
    map_x = Inches(0.5)
    map_y = Inches(1.8)
    map_w = Inches(7.2)
    map_h = Inches(5.0)
    world_map_dots(s, map_x, map_y, map_w, map_h,
                   dot_color=DOT_BG, bg=WHITE)
    # Country highlight circles (노란색 원형, 참고 이미지 스타일)
    # USA (large)
    highlight_country(s, map_x, map_y, map_w, map_h,
                      rx=0.135, ry=0.34, size=1.1,
                      color=YELLOW_HL, label="UNITED\nSTATES",
                      label_color=TEXT_BLACK, label_size=10)
    # Mexico
    highlight_country(s, map_x, map_y, map_w, map_h,
                      rx=0.165, ry=0.50, size=0.55,
                      color=YELLOW_HL, label="MEXICO",
                      label_size=9)
    # Brazil
    highlight_country(s, map_x, map_y, map_w, map_h,
                      rx=0.235, ry=0.72, size=0.7,
                      color=YELLOW_HL, label="BRAZIL",
                      label_size=9)
    # Canada (small)
    highlight_country(s, map_x, map_y, map_w, map_h,
                      rx=0.14, ry=0.20, size=0.4,
                      color=YELLOW_HL, label="CANADA",
                      label_size=8)

    # RIGHT — Country breakdown table
    eyebrow(s, Inches(8.0), Inches(1.85),
            "2025 · AMERICAS BREAKDOWN", color=TEAL)
    data = [
        ["지역", "국가",      "판매량",   "비중"],
        ["북미", "USA",       "약 90.2만 대", "21.8%"],
        ["북미", "Canada",    "약 14.1만 대", "3.4%"],
        ["북미", "Mexico",    "약 5.1만 대",  "1.2%"],
        ["남미", "Brazil (HMB)", "약 18.5만 대", "4.5%"],
        ["기타", "기타",      "약 6.1만 대",  "1.5%"],
        ["합계", "Americas",  "약 134만 대",  "32.4%"],
    ]
    toss_table(s, Inches(8.0), Inches(2.2), Inches(4.8), Inches(3.9),
               data, col_widths=[0.8, 1.5, 1.4, 0.8],
               header_fill=BG_TINT,
               body_size=11, header_size=10,
               highlight_row=4,
               first_col_bold=True,
               center_cols=[2, 3])

    # Bottom insight
    rounded(s, Inches(8.0), Inches(6.25), Inches(4.8), Inches(0.75),
            fill=TEAL_SOFT, radius=0.08)
    txt(s, Inches(8.2), Inches(6.35), Inches(4.5), Inches(0.3),
        "💡 Brazil HMB = 18.5만 대",
        size=11, bold=True, color=TEAL_DEEP)
    txt(s, Inches(8.2), Inches(6.65), Inches(4.5), Inches(0.3),
        "남미는 Tenneco가 이미 점유 확장 중인 시장",
        size=10, color=TEXT_GRAY)
    footer(s, 4)


# -------- Slide 5 — vs GM & Toyota --------
# 말풍선: "이런 식의 그림이 사용되어..." → pie/donut chart
def slide_05_vs_competitors():
    s = add_slide(WHITE)
    slide_title(s, "가장 작지만, 가장 빠르다",
                "Hyundai vs. GM · Toyota — America 2025")

    # LEFT — Donut chart (말풍선 요청 "파이차트 스타일")
    eyebrow(s, Inches(0.6), Inches(1.9),
            "AMERICAS VOLUME SHARE", color=TEAL)
    donut_data = [
        ("GM (320~340만)", 330),
        ("Toyota (280~300만)", 290),
        ("Hyundai (134만)", 134),
    ]
    add_donut_chart(s, Inches(0.4), Inches(2.15),
                    Inches(5.5), Inches(4.5), donut_data)

    # RIGHT — 3 cards
    eyebrow(s, Inches(6.5), Inches(1.9),
            "MARKET POSITION", color=TEAL)
    cards = [
        ("GM", "320–340만 대", "북미 1위",
         "대형 픽업트럭 · SUV 중심의\n압도적 시장 지배력",
         False),
        ("Toyota", "280–300만 대", "하이브리드 절대 강자",
         "북미·남미 모두에서\n높은 신뢰도 보유",
         False),
        ("Hyundai", "134만 대", "가장 빠른 성장세",
         "전기차 · SUV 라인업 강화\n공격적 추격 중",
         True),
    ]
    y0 = Inches(2.2)
    card_h = Inches(1.55)
    card_w = Inches(6.3)
    for i, (brand, volume, pos, desc, highlight) in enumerate(cards):
        y = y0 + (card_h + Inches(0.12)) * i
        fill_c = TEAL_SOFT if highlight else WHITE
        rounded(s, Inches(6.5), y, card_w, card_h,
                fill=fill_c, line=BORDER_GRAY, radius=0.05)
        if highlight:
            rect(s, Inches(6.5), y, Inches(0.08), card_h, fill=TEAL)
        # Brand
        txt(s, Inches(6.7), y + Inches(0.15), Inches(2), Inches(0.35),
            brand, size=16, bold=True,
            color=TEAL if highlight else TEXT_BLACK, font=FONT_EN)
        # Volume
        txt(s, Inches(6.7), y + Inches(0.5), Inches(3.5), Inches(0.5),
            volume, size=26, bold=True,
            color=TEAL_DEEP if highlight else TEXT_BLACK,
            font=FONT_EN)
        # Position
        txt(s, Inches(10.5), y + Inches(0.2), Inches(2.3), Inches(0.3),
            pos, size=11, bold=True, color=TEXT_GRAY,
            align=PP_ALIGN.RIGHT)
        # Desc
        txt(s, Inches(10.5), y + Inches(0.55), Inches(2.3), Inches(0.9),
            desc, size=10, color=TEXT_GRAY,
            align=PP_ALIGN.RIGHT)

    # Bottom insight
    txt(s, Inches(6.5), Inches(7.0), Inches(6.3), Inches(0.3),
        "📈  Volume이 아니라 성장률에 주목하라.",
        size=12, bold=True, color=TEAL_DEEP)
    footer(s, 5)


# -------- Slide 6 — Tenneco Business Ratio in Hyundai --------
def slide_06_tenneco_ratio():
    s = add_slide(WHITE)
    slide_title(s, "현대 섀시 시장의 95%는 경쟁사가 쥐고 있다",
                "Tenneco Business Ratio in Hyundai — Current Status")

    # LEFT: Hero stat
    eyebrow(s, Inches(0.6), Inches(2.0),
            "TENNECO + OTHERS", color=RED)
    txt(s, Inches(0.6), Inches(2.25), Inches(5), Inches(2.0),
        "5%", size=180, bold=True, color=RED, font=FONT_EN)
    txt(s, Inches(0.6), Inches(4.6), Inches(5), Inches(0.4),
        "Tenneco가 점유 중인 비중",
        size=13, color=TEXT_GRAY)

    # RIGHT: Stacked bar + breakdown
    eyebrow(s, Inches(6.2), Inches(2.0),
            "MARKET SHARE BREAKDOWN", color=TEAL)

    # Legend above bar
    legend_y = Inches(2.35)
    legends = [
        ("Mando",    "65%", TEAL),
        ("ZF",       "30%", TEAL_MID),
        ("Tenneco+", "5%",  RED),
    ]
    lx = Inches(6.2)
    for name, pct, c in legends:
        oval(s, lx, legend_y + Inches(0.05), Inches(0.14),
             Inches(0.14), fill=c)
        txt(s, lx + Inches(0.2), legend_y, Inches(1.5), Inches(0.25),
            f"{name} {pct}", size=11, bold=True, color=TEXT_BLACK)
        lx = lx + Inches(1.85)

    # Stacked horizontal bar
    stat_stackbar(s, Inches(6.2), Inches(2.9), Inches(6.6),
                  Inches(0.7),
                  [(65, TEAL, "Mando"),
                   (30, TEAL_MID, "ZF"),
                   (5,  RED, "Tenneco+")])
    # Labels on bar
    txt(s, Inches(6.2), Inches(3.15), Inches(4.29), Inches(0.4),
        "Mando  ·  65%", size=13, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.49), Inches(3.15), Inches(1.98), Inches(0.4),
        "ZF  ·  30%", size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 5% too small to label inside

    # Insight cards (3)
    y_cards = Inches(4.0)
    insights = [
        ("Mando",
         "Volume Dominant",
         "단가 경쟁력 + 국내 생산 허브\n전 차종 대응 인프라"),
        ("ZF",
         "Premium Segment Leader",
         "Air Suspension · 고성능\nHyundai 프리미엄화 전략에 합치"),
        ("New Entrants",
         "Bilstein · S&T",
         "Bilstein: GV90 플래그십 EV\nS&T: Compact car 진입"),
    ]
    card_w = Inches(2.2)
    for i, (h1, h2, body) in enumerate(insights):
        x = Inches(6.2) + (card_w + Inches(0.1)) * i
        rounded(s, x, y_cards, card_w, Inches(2.5),
                fill=BG_TINT, radius=0.08)
        txt(s, x + Inches(0.18), y_cards + Inches(0.15),
            card_w - Inches(0.3), Inches(0.3),
            h1, size=11, bold=True, color=TEAL, font=FONT_EN)
        txt(s, x + Inches(0.18), y_cards + Inches(0.45),
            card_w - Inches(0.3), Inches(0.6),
            h2, size=13, bold=True, color=TEXT_BLACK)
        txt(s, x + Inches(0.18), y_cards + Inches(1.1),
            card_w - Inches(0.3), Inches(1.3),
            body, size=10, color=TEXT_GRAY, spacing=1.3)

    # Bottom insight
    rounded(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.4),
            fill=TEAL_SOFT, radius=0.3)
    txt(s, Inches(0.8), Inches(6.75), Inches(12), Inches(0.4),
        "💡  Mando가 볼륨을, ZF가 프리미엄을 지배 — Tenneco 포지셔닝 재정의가 필요하다.",
        size=12, bold=True, color=TEAL_DEEP,
        anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 6)



# -------- Slide 7 — 4-1. Supply Program Overview --------
# 말풍선: "국기 + 차량사진 + 프로젝트명" 카드 + 세계지도 마커
def slide_07_supply_overview():
    s = add_slide(WHITE)
    slide_title(s, "우리는 이미 9개 프로그램에서 현대와 함께한다",
                "4-1.  테네코에서의 현대 공급 프로그램  ·  CN · IN · BR")

    # Row 1: CHINA (4 cards)
    y_row1 = Inches(1.75)
    eyebrow(s, Inches(0.6), y_row1,
            "CHINA  ·  4 projects", color=TEAL)
    china = [
        ("CN", "CHINA", "GC",    "IX25",          "2014"),
        ("CN", "CHINA", "YC FL", "Verna",         "2019"),
        ("CN", "CHINA", "ADC",   "Elantra 6th",   "2016"),
        ("CN", "CHINA", "CN7c",  "Elantra 7th",   "2020"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.7)
    y_cards1 = y_row1 + Inches(0.35)
    for i, (code, country, proj, veh, yr) in enumerate(china):
        x = Inches(0.6) + (card_w + Inches(0.17)) * i
        vehicle_card(s, x, y_cards1, card_w, card_h,
                     code, country, proj, veh, yr)

    # Row 2: INDIA (2 cards) + BRAZIL (3 cards)
    y_row2 = Inches(4.0)
    eyebrow(s, Inches(0.6), y_row2,
            "INDIA  ·  2 projects", color=TEAL)
    india = [
        ("IN", "INDIA", "QXI", "Venue",      "2019"),
        ("IN", "INDIA", "BI3", "Elite i20",  "2020"),
    ]
    y_cards2 = y_row2 + Inches(0.35)
    for i, (code, country, proj, veh, yr) in enumerate(india):
        x = Inches(0.6) + (card_w + Inches(0.17)) * i
        vehicle_card(s, x, y_cards2, card_w, card_h,
                     code, country, proj, veh, yr)

    eyebrow(s, Inches(6.9), y_row2,
            "BRAZIL  ·  3 projects", color=TEAL)
    brazil = [
        ("BR", "BRAZIL", "BR2",  "HB20",   "2022", False),
        ("BR", "BRAZIL", "SU2b", "Creta",  "2023", False),
        ("BR", "BRAZIL", "BC4b", "New",    "2026", True),
    ]
    card_w2 = Inches(1.95)
    for i, (code, country, proj, veh, yr, new) in enumerate(brazil):
        x = Inches(6.9) + (card_w2 + Inches(0.12)) * i
        vehicle_card(s, x, y_cards2, card_w2, card_h,
                     code, country, proj, veh, yr,
                     new_badge=new, highlight=new)

    # Bottom insight band
    rounded(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.7),
            fill=TEAL_SOFT, radius=0.08)
    txt(s, Inches(0.85), Inches(6.5), Inches(12), Inches(0.25),
        "GLOBAL FOOTPRINT",
        size=10, bold=True, color=TEAL, font=FONT_EN)
    txt(s, Inches(0.85), Inches(6.75), Inches(12), Inches(0.32),
        "총 9개 프로그램 · 3대륙 동시 대응 체제 — BC4b (브라질, 2026) 신규 양산 예정",
        size=12, bold=True, color=TEAL_DEEP)
    footer(s, 7)


# -------- Slide 8 — 4-2. Program Vehicle Volume (Ranked Infographic) --------
# 말풍선: "Sales by country infographics" 스타일 → 랭킹 + 세계지도 번호 마커
def slide_08_volume_infographic():
    s = add_slide(WHITE)
    slide_title(s, "5년간 누적 146만 대의 공급 실적",
                "4-2.  Program Vehicle Volume  ·  Sales by Country (2023~2027)")

    # LEFT: Ranked country cards (Sales by country infographics style)
    eyebrow(s, Inches(0.6), Inches(1.85),
            "CUMULATIVE VOLUME · 5Y TOTAL", color=TEAL)

    # Rank 01: Brazil
    y1 = Inches(2.15)
    rounded(s, Inches(0.6), y1, Inches(6.0), Inches(1.55),
            fill=WHITE, line=BORDER_GRAY, radius=0.05)
    # number badge
    rounded(s, Inches(0.8), y1 + Inches(0.2), Inches(0.7), Inches(0.7),
            fill=TEAL, radius=0.5)
    txt(s, Inches(0.8), y1 + Inches(0.2), Inches(0.7), Inches(0.7),
        "01", size=22, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # flag
    flag(s, Inches(1.7), y1 + Inches(0.25), Inches(0.38), Inches(0.24),
         "BR")
    txt(s, Inches(2.15), y1 + Inches(0.2), Inches(2), Inches(0.3),
        "BRAZIL", size=14, bold=True, color=TEXT_BLACK, font=FONT_EN)
    # Big number
    txt(s, Inches(1.7), y1 + Inches(0.55), Inches(4), Inches(0.7),
        "899,400", size=38, bold=True, color=TEAL_DEEP, font=FONT_EN)
    txt(s, Inches(1.7), y1 + Inches(1.18), Inches(4), Inches(0.25),
        "5-year cumulative · BR2 · SU2b · BC4b",
        size=10, color=TEXT_GRAY, font=FONT_EN)
    # share bar on right
    txt(s, Inches(5.6), y1 + Inches(0.4), Inches(0.9), Inches(0.25),
        "61.5%", size=16, bold=True, color=TEAL,
        align=PP_ALIGN.RIGHT, font=FONT_EN)
    txt(s, Inches(5.6), y1 + Inches(0.7), Inches(0.9), Inches(0.22),
        "share", size=9, color=TEXT_GRAY,
        align=PP_ALIGN.RIGHT, font=FONT_EN)

    # Rank 02: India
    y2 = Inches(3.85)
    rounded(s, Inches(0.6), y2, Inches(6.0), Inches(1.55),
            fill=WHITE, line=BORDER_GRAY, radius=0.05)
    rounded(s, Inches(0.8), y2 + Inches(0.2), Inches(0.7), Inches(0.7),
            fill=TEAL_MID, radius=0.5)
    txt(s, Inches(0.8), y2 + Inches(0.2), Inches(0.7), Inches(0.7),
        "02", size=22, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    flag(s, Inches(1.7), y2 + Inches(0.25), Inches(0.38), Inches(0.24),
         "IN")
    txt(s, Inches(2.15), y2 + Inches(0.2), Inches(2), Inches(0.3),
        "INDIA", size=14, bold=True, color=TEXT_BLACK, font=FONT_EN)
    txt(s, Inches(1.7), y2 + Inches(0.55), Inches(4), Inches(0.7),
        "562,140", size=38, bold=True, color=TEAL_MID, font=FONT_EN)
    txt(s, Inches(1.7), y2 + Inches(1.18), Inches(4), Inches(0.25),
        "5-year cumulative · BI3 · Qxi",
        size=10, color=TEXT_GRAY, font=FONT_EN)
    txt(s, Inches(5.6), y2 + Inches(0.4), Inches(0.9), Inches(0.25),
        "38.5%", size=16, bold=True, color=TEAL_MID,
        align=PP_ALIGN.RIGHT, font=FONT_EN)
    txt(s, Inches(5.6), y2 + Inches(0.7), Inches(0.9), Inches(0.22),
        "share", size=9, color=TEXT_GRAY,
        align=PP_ALIGN.RIGHT, font=FONT_EN)

    # TOTAL
    y3 = Inches(5.55)
    rounded(s, Inches(0.6), y3, Inches(6.0), Inches(0.65),
            fill=TEAL_SOFT, radius=0.05)
    txt(s, Inches(0.9), y3 + Inches(0.15), Inches(2), Inches(0.38),
        "TOTAL", size=11, bold=True, color=TEAL_DEEP, font=FONT_EN)
    txt(s, Inches(2.8), y3 + Inches(0.05), Inches(3), Inches(0.55),
        "1,461,540", size=28, bold=True, color=TEAL_DEEP, font=FONT_EN)
    txt(s, Inches(5.6), y3 + Inches(0.18), Inches(0.9), Inches(0.35),
        "100%", size=14, bold=True, color=TEAL_DEEP,
        align=PP_ALIGN.RIGHT, font=FONT_EN)

    # RIGHT: World map with numbered markers
    map_x = Inches(7.0)
    map_y = Inches(2.15)
    map_w = Inches(5.9)
    map_h = Inches(4.05)
    world_map_dots(s, map_x, map_y, map_w, map_h)
    # Brazil marker #01
    br_cx = map_x + Emu(int(map_w * 0.24))
    br_cy = map_y + Emu(int(map_h * 0.72))
    numbered_marker(s, br_cx, br_cy, "01", color=TEAL, size=0.55)
    txt(s, br_cx - Inches(0.8), br_cy + Inches(0.32),
        Inches(1.6), Inches(0.25),
        "BRAZIL · 899K", size=10, bold=True, color=TEXT_BLACK,
        align=PP_ALIGN.CENTER, font=FONT_EN)
    # India marker #02
    in_cx = map_x + Emu(int(map_w * 0.72))
    in_cy = map_y + Emu(int(map_h * 0.48))
    numbered_marker(s, in_cx, in_cy, "02", color=TEAL_MID, size=0.5)
    txt(s, in_cx - Inches(0.8), in_cy + Inches(0.3),
        Inches(1.6), Inches(0.25),
        "INDIA · 562K", size=10, bold=True, color=TEXT_BLACK,
        align=PP_ALIGN.CENTER, font=FONT_EN)

    # Year-by-year mini table at bottom
    y_tbl = Inches(6.4)
    data_y = [
        ["",       "2023",    "2024",    "2025",    "2026",    "2027"],
        ["India",  "219,900", "196,820", "138,560", "6,860",   "–"],
        ["Brazil", "137,600", "148,800", "216,000", "216,000", "181,000"],
    ]
    toss_table(s, Inches(0.6), y_tbl, Inches(12.2), Inches(0.65),
               data_y, col_widths=[1.2, 1.5, 1.5, 1.5, 1.5, 1.5],
               body_size=10, header_size=9,
               center_cols=[1, 2, 3, 4, 5],
               header_fill=BG_TINT)
    footer(s, 8)


# -------- Slide 9 — 4-3. Market Status by Region --------
def slide_09_market_status():
    s = add_slide(WHITE)
    slide_title(s, "Brazil만 이기고 있다 — 이유를 복제하라",
                "4-3.  Market Status by Region  ·  BHMC · HMIL · HMB")

    # Table: 3 rows x 4 cols (Region / Customer / Status / Root Cause)
    y0 = Inches(1.95)
    row_h = Inches(1.55)
    col_widths = [1.8, 3.2, 2.1, 6.0]

    # Header
    header_y = y0
    hx = [Inches(0.6)]
    total_w = sum(col_widths)
    cum = 0
    for cw in col_widths:
        cum += cw
    xs = [Inches(0.6)]
    acc = 0.6
    for cw in col_widths:
        acc += cw
        xs.append(Inches(acc))
    # Header row
    for i, head in enumerate(["REGION", "CUSTOMER", "STATUS", "ROOT CAUSE"]):
        txt(s, xs[i], header_y, Inches(col_widths[i]), Inches(0.35),
            head, size=10, bold=True, color=TEXT_GRAY, font=FONT_EN)
    hairline(s, Inches(0.6), header_y + Inches(0.42),
             Inches(sum(col_widths)), color=TEXT_BLACK, weight=1.0)

    # Rows
    rows = [
        ("Tenneco\nChina",
         "BHMC",
         "(Beijing Hyundai\nMotor Company)",
         ("Sales Decline", "red"),
         ["ICE → EV 전환 적응 부족",
          "HMC China 전체 매출 축소 → 신규 사업 정체"]),
        ("Tenneco\nIndia",
         "HMIL",
         "(Hyundai Motor\nIndia Limited)",
         ("Sales Decline", "red"),
         ["Post-Pandemic Order Gap",
          "CA/MRS 미흡 · After-sales 지원 부족",
          "신규 차종 수주 실패"]),
        ("Tenneco\nBrazil",
         "HMB",
         "(Hyundai Motor\nBrazil)",
         ("Market Share ↑", "green"),
         ["HMC Brazil 생산량 증가 + 구매팀 전략적 소싱",
          "Mando 현지 생산 부재 → Tenneco 전략적 우위"]),
    ]
    row_y = y0 + Inches(0.55)
    for i, (region, cust_k, cust_en, (status, tone), causes) in enumerate(rows):
        y = row_y + row_h * i
        # background for Brazil row (highlight)
        if i == 2:
            rect(s, Inches(0.6), y - Inches(0.1),
                 Inches(sum(col_widths)), row_h,
                 fill=TEAL_SOFT)
        # Region
        txt(s, xs[0], y, Inches(col_widths[0]), Inches(1.0),
            region, size=13, bold=True, color=TEXT_BLACK, spacing=1.1)
        # Customer
        txt(s, xs[1], y, Inches(col_widths[1]), Inches(0.35),
            cust_k, size=14, bold=True, color=TEXT_BLACK, font=FONT_EN)
        txt(s, xs[1], y + Inches(0.4), Inches(col_widths[1]),
            Inches(0.6),
            cust_en, size=10, color=TEXT_GRAY, font=FONT_EN,
            spacing=1.2)
        # Status pill
        pill_w = Inches(1.7)
        status_pill(s, xs[2], y + Inches(0.25),
                    pill_w, Inches(0.35), status, tone=tone)
        # Root cause bullets
        for j, c in enumerate(causes):
            yy = y + Inches(0.1) + Inches(0.35 * j)
            # bullet dot
            oval(s, xs[3] + Inches(0.05), yy + Inches(0.11),
                 Inches(0.1), Inches(0.1),
                 fill=(GREEN if i == 2 else RED))
            txt(s, xs[3] + Inches(0.25), yy,
                Inches(col_widths[3] - 0.3), Inches(0.35),
                c, size=11, color=TEXT_BLACK)
        # row separator
        if i < 2:
            hairline(s, Inches(0.6), y + row_h - Inches(0.1),
                     Inches(sum(col_widths)))

    # Bottom insight
    rounded(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.45),
            fill=BG_TINT, radius=0.3)
    txt(s, Inches(0.85), Inches(6.7), Inches(12), Inches(0.45),
        "💡  중국·인도에서 잃은 만큼, 브라질의 성공 패턴(구매팀 전략 소싱 + 현지 공급 우위)을 수평 확산해야 한다.",
        size=11, bold=True, color=TEAL_DEEP,
        anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 9)



# -------- Slide 10 — 5. Competitor Overview (with NA & SA plant maps) --------
# 말풍선: "북미·남미 지역의 ZF, 만도 공장(심플하게 이름만 표시) 지도"
def slide_10_competitor_overview():
    s = add_slide(WHITE)
    slide_title(s, "4개의 경쟁사, 각기 다른 무기",
                "5.  Tenneco Competitor  ·  ZF · Mando · Bilstein · S&T")

    # LEFT: Competitor profile cards (2x2)
    eyebrow(s, Inches(0.6), Inches(1.9),
            "COMPETITOR PROFILES", color=TEAL)
    comps = [
        ("Mando",    "Korea",   "Volume Dominant",
         "65% 점유 · 국내 생산·R&D\n즉시 대응력"),
        ("ZF",       "Germany", "Premium · Air Susp.",
         "30% · 고성능 라인 인지도\nHMC 프리미엄 전략 합치"),
        ("Bilstein", "Germany", "GV90 Flagship EV 진입",
         "High-end Air Shock+ECS\n소량·프리미엄"),
        ("S&T",      "Korea",   "Compact 진입",
         "국내 2nd Tier\n가격 경쟁력 보유"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(2.1)
    y_base = Inches(2.2)
    for i, (name, country, tag, desc) in enumerate(comps):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (card_w + Inches(0.15)) * col
        y = y_base + (card_h + Inches(0.15)) * row
        rounded(s, x, y, card_w, card_h,
                fill=WHITE, line=BORDER_GRAY, radius=0.05)
        # Name
        txt(s, x + Inches(0.2), y + Inches(0.15),
            card_w - Inches(0.4), Inches(0.4),
            name, size=20, bold=True, color=TEXT_BLACK, font=FONT_EN)
        # Country tag
        rounded(s, x + card_w - Inches(1.0),
                y + Inches(0.22),
                Inches(0.8), Inches(0.28),
                fill=BG_TINT, radius=0.5)
        txt(s, x + card_w - Inches(1.0), y + Inches(0.22),
            Inches(0.8), Inches(0.28),
            country, size=9, bold=True, color=TEXT_GRAY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        # Tag line
        txt(s, x + Inches(0.2), y + Inches(0.65),
            card_w - Inches(0.4), Inches(0.3),
            tag, size=12, bold=True, color=TEAL)
        # Description
        txt(s, x + Inches(0.2), y + Inches(1.05),
            card_w - Inches(0.4), card_h - Inches(1.15),
            desc, size=10, color=TEXT_GRAY, spacing=1.3)

    # RIGHT: North America plant map (simplified — 말풍선 요청 "심플하게 이름만")
    map_na_x = Inches(6.9)
    map_na_y = Inches(2.2)
    map_na_w = Inches(6.1)
    map_na_h = Inches(2.5)
    eyebrow(s, map_na_x, map_na_y - Inches(0.3),
            "NORTH AMERICA PLANTS", color=TEAL)
    # simple NA outline
    rounded(s, map_na_x, map_na_y, map_na_w, map_na_h,
            fill=BG_TINT, radius=0.03)
    # labels only (state-name style, similar to Korean news graphic)
    na_points = [
        ("Detroit, MI",   "ZF",     0.45, 0.38),
        ("Gainesville, GA","ZF",    0.55, 0.70),
        ("Opelika, AL",   "Mando",  0.50, 0.75),
        ("Monterrey, MX", "Mando",  0.35, 0.92),
        ("Arteaga, MX",   "ZF",     0.30, 0.95),
    ]
    # NOTE: the above 0.95 may go outside; clamp visually
    for label, brand, rx, ry in na_points:
        cx = map_na_x + Emu(int(map_na_w * rx))
        cy = map_na_y + Emu(int(map_na_h * min(ry, 0.85)))
        # pin
        oval(s, cx - Inches(0.06), cy - Inches(0.06),
             Inches(0.12), Inches(0.12),
             fill=(TEAL_DEEP if brand == "ZF" else TEAL))
        # label
        txt(s, cx - Inches(0.9), cy + Inches(0.1),
            Inches(1.8), Inches(0.22),
            f"{label} · {brand}", size=8, bold=True,
            color=TEXT_BLACK,
            align=PP_ALIGN.CENTER, font=FONT_EN)

    # RIGHT BOTTOM: South America plant map
    map_sa_y = Inches(4.9)
    map_sa_h = Inches(2.0)
    eyebrow(s, map_na_x, map_sa_y - Inches(0.3),
            "SOUTH AMERICA PLANTS", color=TEAL)
    rounded(s, map_na_x, map_sa_y, map_na_w, map_sa_h,
            fill=BG_TINT, radius=0.03)
    sa_points = [
        ("São Bernardo",    "ZF",    0.55, 0.30),
        ("Sorocaba",        "ZF",    0.52, 0.55),
        ("Limeira",         "Mando", 0.48, 0.45),
        ("San Francisco, AR","ZF",   0.30, 0.78),
    ]
    for label, brand, rx, ry in sa_points:
        cx = map_na_x + Emu(int(map_na_w * rx))
        cy = map_sa_y + Emu(int(map_sa_h * ry))
        oval(s, cx - Inches(0.06), cy - Inches(0.06),
             Inches(0.12), Inches(0.12),
             fill=(TEAL_DEEP if brand == "ZF" else TEAL))
        txt(s, cx - Inches(0.9), cy + Inches(0.1),
            Inches(1.8), Inches(0.22),
            f"{label} · {brand}", size=8, bold=True,
            color=TEXT_BLACK,
            align=PP_ALIGN.CENTER, font=FONT_EN)
    footer(s, 10)


# -------- Slide 11 — 5-1. Competitor의 특수성 --------
def slide_11_competitor_specificity():
    s = add_slide(WHITE)
    slide_title(s, "그들이 이기는 것은 기술이 아닌 '속도와 근접성'이다",
                "5-1.  Tenneco Competitor의 특수성  ·  Market Status + Gap Analysis")

    # TOP: Market Status (compact)
    eyebrow(s, Inches(0.6), Inches(1.85),
            "MARKET STATUS & BUSINESS BALANCE", color=TEAL)
    # Stack bar
    stat_stackbar(s, Inches(0.6), Inches(2.2), Inches(12.2),
                  Inches(0.45),
                  [(65, TEAL, "Mando"),
                   (30, TEAL_MID, "ZF"),
                   (5,  RED, "Others")])
    # Labels
    txt(s, Inches(0.6), Inches(2.22), Inches(7.93), Inches(0.41),
        "Mando  ·  65%", size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.53), Inches(2.22), Inches(3.66), Inches(0.41),
        "ZF  ·  30%", size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(12.2), Inches(2.22), Inches(0.6), Inches(0.41),
        "5%", size=9, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Recent entrants callout
    txt(s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.3),
        "최근 진입:  Bilstein (GV90 Flagship EV)  ·  S&T (Compact Car)",
        size=11, color=TEXT_GRAY)

    # Divider
    hairline(s, Inches(0.6), Inches(3.3), Inches(12.2))

    # BOTTOM: 3 Key Success Factors (Gap Analysis)
    eyebrow(s, Inches(0.6), Inches(3.5),
            "COMPETITORS' KEY SUCCESS FACTORS  ·  GAP ANALYSIS", color=TEAL)
    factors = [
        ("01", "Real-time Response\nInfrastructure",
         "경쟁사는 한국 내 R&D + Plant 운영.\nHyundai 엔지니어 요청에 즉시 대응 가능.",
         "Tenneco: 국내 생산 허브 부재 + 지원 인력 제한.\n→ 물리적 대응 속도 한계"),
        ("02", "Localized Business\nAlignment",
         "경쟁사는 한국인 전담팀 운영.\nHyundai 구매·개발 프로세스에 완전 동기화.",
         "Tenneco: 해외 Plant/R&D 조율 → 커뮤니케이션 지연 발생 (약점)"),
        ("03", "Extreme Quality\nStandards",
         "Hyundai: 반도체 생산급 Cleanroom (KYB 기준) 요구.\nMando·ZF는 국내 전용 라인으로 충족.",
         "Tenneco: 국내 전용 라인 부재 → 충족 전략 필요"),
    ]
    col_w = Inches(4.0)
    y_f = Inches(3.85)
    f_h = Inches(3.1)
    for i, (num, title, comp_has, gap) in enumerate(factors):
        x = Inches(0.6) + (col_w + Inches(0.1)) * i
        rounded(s, x, y_f, col_w, f_h,
                fill=BG_TINT, radius=0.06)
        # number
        txt(s, x + Inches(0.2), y_f + Inches(0.15),
            col_w - Inches(0.3), Inches(0.6),
            num, size=38, bold=True, color=TEAL, font=FONT_EN)
        # title
        txt(s, x + Inches(0.2), y_f + Inches(0.7),
            col_w - Inches(0.3), Inches(0.7),
            title, size=14, bold=True, color=TEXT_BLACK, spacing=1.15)
        # Competitors have
        txt(s, x + Inches(0.2), y_f + Inches(1.45),
            col_w - Inches(0.3), Inches(0.22),
            "COMPETITORS",
            size=9, bold=True, color=TEAL, font=FONT_EN)
        txt(s, x + Inches(0.2), y_f + Inches(1.65),
            col_w - Inches(0.3), Inches(0.7),
            comp_has, size=10, color=TEXT_BLACK, spacing=1.3)
        # Tenneco gap
        txt(s, x + Inches(0.2), y_f + Inches(2.3),
            col_w - Inches(0.3), Inches(0.22),
            "TENNECO GAP",
            size=9, bold=True, color=RED, font=FONT_EN)
        txt(s, x + Inches(0.2), y_f + Inches(2.5),
            col_w - Inches(0.3), Inches(0.55),
            gap, size=10, color=TEXT_GRAY, spacing=1.3)
    footer(s, 11)


# -------- Slide 12 — 5-2. Competitor Matrix --------
def slide_12_competitor_matrix():
    s = add_slide(WHITE)
    slide_title(s, "한눈에 본 4사 경쟁 포지션",
                "5-2.  Competitor Matrix  ·  Tenneco만 Plant · R&D 국내 부재")

    # Supplier matrix table
    data = [
        ["Supplier", "Office\n(KR)", "Plant\n(KR)", "R&D\n(KR)",
         "Strength", "Weakness"],
        ["Mando",
         "O", "O", "O",
         "빠른 대응 · 가격 경쟁력",
         "2-valve 등 신기술 부재"],
        ["ZF",
         "O", "O", "O",
         "고성능 인지도 · 빠른 개발",
         "가격 高 · Air+Shock 조합 부재"],
        ["Bilstein",
         "O", "X", "O",
         "고성능 프리미엄 인지도",
         "가격 高 · 2-valve MacPherson 부재"],
        ["Tenneco",
         "O", "X", "X",
         "HMC 평가 — 가격 · 성능 합리적",
         "HMC 고성능 개발 경험 부재"],
    ]
    col_widths = [2.0, 1.3, 1.3, 1.3, 3.5, 3.1]
    tbl = toss_table(s, Inches(0.6), Inches(1.95),
                     Inches(12.2), Inches(4.2),
                     data, col_widths=col_widths,
                     header_fill=BG_TINT,
                     body_size=11, header_size=10,
                     highlight_row=4,
                     first_col_bold=True,
                     center_cols=[1, 2, 3])
    # Manually color O/X symbols via re-running paragraph runs
    for ri in range(1, 5):
        for ci in [1, 2, 3]:
            cell = tbl.cell(ri, ci)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    val = r.text.strip()
                    r.font.size = Pt(16)
                    r.font.bold = True
                    if val == "O":
                        r.font.color.rgb = TEAL
                    elif val == "X":
                        r.font.color.rgb = RED

    # Bottom insight
    rounded(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.7),
            fill=TEAL_SOFT, radius=0.06)
    txt(s, Inches(0.85), Inches(6.5), Inches(12), Inches(0.25),
        "POSITIONING",
        size=10, bold=True, color=TEAL, font=FONT_EN)
    txt(s, Inches(0.85), Inches(6.78), Inches(12), Inches(0.3),
        "Plant · R&D 국내 부재를 극복하는 유일한 길은 '기술 차별성 + 가격 합리성 + 유럽 트랙 레코드'.",
        size=12, bold=True, color=TEAL_DEEP)
    footer(s, 12)



# -------- Slide 13 — 6. 현대 공략 방법 --------
# 말풍선 요청: "미국과 브라질만 부각되어 보이도록" — USA·Brazil 하이라이트 지도
def slide_13_attack_strategy():
    s = add_slide(WHITE)
    slide_title(s, "기술은 있다. 가격과 공장을 맞추면 이긴다",
                "6.  현대 공략 방법  ·  N.America Opportunity + S.America Response")

    # LEFT: Americas map with USA + Brazil highlighted (말풍선 요청)
    map_x = Inches(0.6)
    map_y = Inches(1.9)
    map_w = Inches(4.2)
    map_h = Inches(5.2)
    # background
    rect(s, map_x, map_y, map_w, map_h, fill=BG_TINT)

    # Simplified Americas continents
    # North America (USA) — HIGHLIGHTED (filled with brand blue)
    # Canada — grayed out
    rounded(s, map_x + Inches(0.3), map_y + Inches(0.25),
            Inches(3.4), Inches(1.0),
            fill=BG_GRAY, radius=0.15)
    txt(s, map_x + Inches(0.3), map_y + Inches(0.55),
        Inches(3.4), Inches(0.3),
        "CANADA", size=10, bold=True, color=TEXT_LIGHT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

    # USA — HIGHLIGHTED
    rounded(s, map_x + Inches(0.45), map_y + Inches(1.35),
            Inches(3.0), Inches(1.0),
            fill=TEAL, radius=0.12, line=TEAL_DEEP, line_w=1.5)
    txt(s, map_x + Inches(0.45), map_y + Inches(1.35),
        Inches(3.0), Inches(1.0),
        "UNITED STATES", size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

    # Mexico — partial highlight (Celaya/Jalisco indication)
    rounded(s, map_x + Inches(0.8), map_y + Inches(2.45),
            Inches(2.2), Inches(0.65),
            fill=BG_GRAY, radius=0.2)
    txt(s, map_x + Inches(0.8), map_y + Inches(2.45),
        Inches(2.2), Inches(0.3),
        "MEXICO", size=10, bold=True, color=TEXT_GRAY,
        align=PP_ALIGN.CENTER, font=FONT_EN)
    # Celaya/Jalisco point
    oval(s, map_x + Inches(1.6), map_y + Inches(2.73),
         Inches(0.18), Inches(0.18), fill=LIME, line=WHITE, line_w=1.2)
    txt(s, map_x + Inches(0.8), map_y + Inches(2.9),
        Inches(2.2), Inches(0.22),
        "Celaya", size=9, bold=True, color=TEAL_DEEP,
        align=PP_ALIGN.CENTER, font=FONT_EN)

    # Connector (thin line between NA block and SA block)
    line(s, map_x + Inches(1.85), map_y + Inches(3.15),
         map_x + Inches(1.85), map_y + Inches(3.7),
         color=BORDER_GRAY, weight=1.0)

    # Brazil — HIGHLIGHTED
    rounded(s, map_x + Inches(1.2), map_y + Inches(3.75),
            Inches(2.4), Inches(1.2),
            fill=TEAL, radius=0.12, line=TEAL_DEEP, line_w=1.5)
    txt(s, map_x + Inches(1.2), map_y + Inches(3.75),
        Inches(2.4), Inches(1.2),
        "BRAZIL", size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

    # Argentina — grayed
    rounded(s, map_x + Inches(0.9), map_y + Inches(5.1),
            Inches(1.4), Inches(0.5),
            fill=BG_GRAY, radius=0.2)
    txt(s, map_x + Inches(0.9), map_y + Inches(5.1),
        Inches(1.4), Inches(0.5),
        "ARGENTINA", size=9, bold=True, color=TEXT_LIGHT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

    # Legend below map
    txt(s, map_x, map_y + map_h - Inches(0.35),
        map_w, Inches(0.3),
        "▍ 진한 파랑 = 공략 우선 지역",
        size=9, color=TEAL, align=PP_ALIGN.CENTER)

    # RIGHT: Two stacked strategy cards
    # TOP — Opportunities (N.America)
    top_x = Inches(5.1)
    top_y = Inches(1.9)
    top_w = Inches(7.8)
    top_h = Inches(2.55)
    rounded(s, top_x, top_y, top_w, top_h,
            fill=WHITE, line=BORDER_GRAY, radius=0.05)
    eyebrow(s, top_x + Inches(0.25), top_y + Inches(0.18),
            "OPPORTUNITIES  ·  N.AMERICA", color=TEAL)
    txt(s, top_x + Inches(0.25), top_y + Inches(0.42),
        top_w - Inches(0.5), Inches(0.5),
        "기술 우위는 이미 입증됐다",
        size=20, bold=True, color=TEXT_BLACK)
    # Block 1
    txt(s, top_x + Inches(0.25), top_y + Inches(1.0),
        top_w - Inches(0.5), Inches(0.3),
        "Technological Advantage",
        size=11, bold=True, color=TEAL)
    txt(s, top_x + Inches(0.25), top_y + Inches(1.3),
        top_w - Inches(0.5), Inches(0.4),
        "HMC 엔지니어링팀이 Tenneco 밸브 기술(Conventional & ECS)의 우월성 인정",
        size=11, color=TEXT_BLACK)
    # Block 2
    txt(s, top_x + Inches(0.25), top_y + Inches(1.8),
        top_w - Inches(0.5), Inches(0.3),
        "Performance Evaluation",
        size=11, bold=True, color=TEAL)
    txt(s, top_x + Inches(0.25), top_y + Inches(2.1),
        top_w - Inches(0.5), Inches(0.4),
        "CVSA2 평가에서 ZF · Mando 대비 성능 우위 확인 (High-Performance Vehicle Team)",
        size=11, color=TEXT_BLACK)

    # BOTTOM — Response Strategy (S.America / New CEO José Muñoz)
    bot_y = Inches(4.6)
    bot_h = Inches(2.55)
    rounded(s, top_x, bot_y, top_w, bot_h,
            fill=TEAL_SOFT, radius=0.05)
    rect(s, top_x, bot_y, Inches(0.08), bot_h, fill=TEAL)
    eyebrow(s, top_x + Inches(0.25), bot_y + Inches(0.18),
            "RESPONSE STRATEGY  ·  S.AMERICA", color=TEAL)
    txt(s, top_x + Inches(0.25), bot_y + Inches(0.42),
        top_w - Inches(0.5), Inches(0.5),
        "가격과 공장을 현지화한다",
        size=20, bold=True, color=TEAL_DEEP)
    # Sub text — New CEO José Muñoz
    txt(s, top_x + Inches(0.25), bot_y + Inches(0.95),
        top_w - Inches(0.5), Inches(0.3),
        "Hyundai New CEO — José Muñoz  ·  Cost Reduction Initiative 대응",
        size=11, italic=True, color=TEXT_GRAY)
    # Two-column action items
    items = [
        ("Strategic Pricing Policy",
         "Mando와 동등한 단가 경쟁력 확보"),
        ("N.America Manufacturing",
         "Celaya(Mexico) Plant를 KYB 수준\n품질·공정 관리로 업그레이드"),
    ]
    for i, (h, body) in enumerate(items):
        x = top_x + Inches(0.25) + Inches(3.85) * i
        w = Inches(3.6)
        txt(s, x, bot_y + Inches(1.4), w, Inches(0.3),
            h, size=11, bold=True, color=TEAL)
        txt(s, x, bot_y + Inches(1.7), w, Inches(0.8),
            body, size=11, color=TEXT_BLACK, spacing=1.3)
    footer(s, 13)


# -------- Slide 14 — 6-1. ECS Opportunity (CVSA2 Program) --------
def slide_14_cvsa2():
    s = add_slide(WHITE)
    slide_title(s, "Ioniq 6 N이 우리의 증명 무대다",
                "6-1.  ECS Opportunity  ·  CVSA2 Program Background")

    # TOP band: Target
    y_top = Inches(1.85)
    rounded(s, Inches(0.6), y_top, Inches(12.2), Inches(1.05),
            fill=TEAL_SOFT, radius=0.06)
    eyebrow(s, Inches(0.8), y_top + Inches(0.12),
            "TARGET VEHICLE", color=TEAL)
    txt(s, Inches(0.8), y_top + Inches(0.4),
        Inches(3.5), Inches(0.55),
        "Ioniq 6 N", size=28, bold=True, color=TEAL_DEEP,
        font=FONT_EN)
    # Competitor chips
    chips = [
        ("Base",   "ZF · one valve",     WHITE, TEXT_BLACK, BORDER_GRAY),
        ("Test 1", "Mando · two valve",  WHITE, TEXT_BLACK, BORDER_GRAY),
        ("Test 2", "Tenneco · CVSA2",    TEAL,  WHITE,      TEAL),
    ]
    x_chip = Inches(4.6)
    for label, body, bg, fg, line_c in chips:
        rounded(s, x_chip, y_top + Inches(0.33),
                Inches(2.65), Inches(0.55),
                fill=bg, line=line_c, radius=0.2)
        txt(s, x_chip + Inches(0.15), y_top + Inches(0.34),
            Inches(1.0), Inches(0.22),
            label, size=8, bold=True, color=(LIME if bg == TEAL else TEAL),
            font=FONT_EN)
        txt(s, x_chip + Inches(0.15), y_top + Inches(0.55),
            Inches(2.4), Inches(0.3),
            body, size=11, bold=True, color=fg, font=FONT_EN)
        x_chip = x_chip + Inches(2.72)

    # MIDDLE: Timeline (5 nodes)
    y_tl = Inches(3.15)
    eyebrow(s, Inches(0.6), y_tl,
            "SCHEDULE  ·  DEMO CAR TIMELINE", color=TEAL)
    # line connector
    line(s, Inches(1.2), y_tl + Inches(1.0),
         Inches(12.2), y_tl + Inches(1.0),
         color=BORDER_GRAY, weight=2.0)
    # nodes
    tl_nodes = [
        ("CW42",        "Oct 13–19",     "샘플 선적"),
        ("CW46",        "Nov 11–16",     "전장 설치"),
        ("CW47–49",     "Nov 17–Dec 7",  "HW/SW 튜닝 #1"),
        ("CW16",        "2026 Apr 13–17", "HW/SW 튜닝 #2"),
        ("Jun 2026",    "End of Q2",     "개선안 수립"),
    ]
    node_w = Inches(2.42)
    for i, (cw, date, desc) in enumerate(tl_nodes):
        cx = Inches(1.2) + Inches(2.72) * i
        # dot
        is_key = i in [2, 4]
        dot_r = Inches(0.14 if is_key else 0.1)
        oval(s, cx - dot_r, y_tl + Inches(1.0) - dot_r,
             dot_r * 2, dot_r * 2,
             fill=(TEAL if is_key else TEAL_MID),
             line=WHITE, line_w=2)
        # CW label
        txt(s, cx - node_w / 2, y_tl + Inches(0.35),
            node_w, Inches(0.3),
            cw, size=13, bold=True, color=TEAL_DEEP,
            align=PP_ALIGN.CENTER, font=FONT_EN)
        # Date
        txt(s, cx - node_w / 2, y_tl + Inches(0.63),
            node_w, Inches(0.25),
            date, size=9, color=TEXT_GRAY,
            align=PP_ALIGN.CENTER, font=FONT_EN)
        # Desc below line
        txt(s, cx - node_w / 2, y_tl + Inches(1.2),
            node_w, Inches(0.3),
            desc, size=10, bold=True, color=TEXT_BLACK,
            align=PP_ALIGN.CENTER)

    # BOTTOM: Tenneco Strength — HMC view (6 points)
    y_btm = Inches(5.1)
    eyebrow(s, Inches(0.6), y_btm,
            "TENNECO STRENGTH  ·  HMC VIEW", color=TEAL)
    strengths = [
        ("2-valve MacPherson + Double Wishbone 대응",
         "차세대 제품 필수 조건 충족",           True),
        ("가격 경쟁력 우위",
         "HMC 내부 평가 — 경쟁사 대비 합리적",   False),
        ("유럽 현지 시승 긍정 피드백",
         "HMC 엔지니어 긍정 반응 확보",         False),
        ("Mando 1-valve 성능 한계 인지",
         "HMC R&D 내부에서 이미 인지",           False),
        ("Demo Car = RFQ 검증",
         "ZF 동등 성능 → Tenneco 우위 확보",   True),
        ("N Brand → Genesis 수직 확산",
         "N Brand 검증 제품은 Genesis로 이전",  False),
    ]
    col_w = Inches(6.05)
    for i, (title, sub, key) in enumerate(strengths):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + Inches(0.1)) * col
        y = y_btm + Inches(0.3) + Inches(0.58) * row
        # Number
        txt(s, x, y, Inches(0.4), Inches(0.3),
            f"0{i+1}", size=14, bold=True,
            color=(RED if key else TEAL),
            font=FONT_EN)
        # Title
        txt(s, x + Inches(0.45), y - Inches(0.02),
            col_w - Inches(0.5), Inches(0.3),
            title, size=11, bold=True,
            color=(RED if key else TEXT_BLACK))
        # Sub
        txt(s, x + Inches(0.45), y + Inches(0.25),
            col_w - Inches(0.5), Inches(0.25),
            sub, size=9, color=TEXT_GRAY)
    footer(s, 14)



# -------- Slide 15 — 6-2. 북미 시장 ECS 공급안 (LX3 / LQ2a 계열) --------
# 말풍선 요청: 복잡한 운영안에서 북미 + ECS 만 추출한 심플 표
def slide_15_ecs_lx3_lq2a():
    s = add_slide(WHITE)
    slide_title(s, "북미 시장 고급옵션 — ECS 공급안 ①",
                "6-2.  LX3 / LQ2a 계열  ·  Palisade · Santa Fe  ·  현재 & 향후")

    # Legend
    lg_y = Inches(1.85)
    legs = [
        ("ECS 적용",  TEAL,  WHITE),
        ("CONV",      BG_GRAY, TEXT_BLACK),
        ("Tenneco 진입 기회", LIME, TEXT_BLACK),
    ]
    lg_x = Inches(0.6)
    for t, c, tc in legs:
        rounded(s, lg_x, lg_y, Inches(1.8), Inches(0.3),
                fill=c, radius=0.3)
        txt(s, lg_x, lg_y, Inches(1.8), Inches(0.3),
            t, size=9, bold=True, color=tc,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        lg_x = lg_x + Inches(1.95)

    # Main table — simplified
    data = [
        ["차종",       "해외명",        "북미 구분",    "현재 ECS 공급사",  "PE 이후 ECS 계획"],
        ["LX3",        "Palisade",      "SELF / XRT PRO", "Mando SFD3 (SELF만)", "Mando 유지 예상"],
        ["LX3 PE",     "Palisade PE",   "SELF · ECS",   "– (신규)",            "ECS 확대 · Tenneco 진입 기회"],
        ["LX3a HEV",   "Palisade HEV",  "SELF",         "Mando SFD3",          "Mando 유지"],
        ["LQ2a",       "Santa Fe",      "XPRO SELF",    "Mando SFD3 (용량증대)","Mando 유지"],
        ["LQ2a PE",    "Santa Fe PE",   "CONV + ECS",   "– (신규)",            "ECS 확대 · Tenneco 진입 기회"],
    ]
    tbl = toss_table(s, Inches(0.6), Inches(2.35),
                     Inches(12.2), Inches(3.4),
                     data,
                     col_widths=[1.5, 2.0, 2.3, 3.0, 3.4],
                     header_fill=TEAL_DEEP,
                     header_color=WHITE,
                     body_size=11, header_size=10,
                     center_cols=[2])
    # Color rows 2 & 5 (PE with Tenneco opportunity)
    for ri in [2, 5]:
        for ci in range(5):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xDC)
        # add colored bar on first cell (emulate by changing first col fill to LIME)
        c0 = tbl.cell(ri, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = LIME

    # Bottom: key attack points (2 callouts)
    y_cta = Inches(5.95)
    ctas = [
        ("🎯",  "LX3 PE · Palisade PE",
         "ECS 신규 옵션 확대 구간. Mando 단독 공급 전환 전 Tenneco 선제 제안 필요."),
        ("🎯",  "LQ2a PE · Santa Fe PE",
         "CONV → ECS 확대 라인. 북미 전용 사양 — Tenneco CVSA2 제안 최적 타이밍."),
    ]
    cta_w = Inches(6.0)
    for i, (icon, title, body) in enumerate(ctas):
        x = Inches(0.6) + (cta_w + Inches(0.2)) * i
        rounded(s, x, y_cta, cta_w, Inches(1.1),
                fill=TEAL_SOFT, radius=0.08)
        rect(s, x, y_cta, Inches(0.08), Inches(1.1), fill=TEAL)
        txt(s, x + Inches(0.2), y_cta + Inches(0.1),
            Inches(0.45), Inches(0.4),
            icon, size=18, color=TEAL)
        txt(s, x + Inches(0.7), y_cta + Inches(0.15),
            cta_w - Inches(0.8), Inches(0.3),
            title, size=12, bold=True, color=TEAL_DEEP)
        txt(s, x + Inches(0.7), y_cta + Inches(0.48),
            cta_w - Inches(0.8), Inches(0.6),
            body, size=10, color=TEXT_BLACK, spacing=1.3)
    footer(s, 15)


# -------- Slide 16 — 6-3. 북미 시장 ECS 공급안 (MV / ME / MV GT 계열) --------
def slide_16_ecs_mv_me():
    s = add_slide(WHITE)
    slide_title(s, "북미 시장 고급옵션 — ECS 공급안 ②",
                "6-3.  MV / ME / MVa / MEa / MV GT 계열  ·  Tucson · Santa Cruz  ·  현재 & 향후")

    # Main table — 북미 ECS만 추출
    data = [
        ["차종",    "해외명",           "북미 구분",       "현재 ECS",                  "PE 이후 ECS 계획"],
        ["MV",      "Tucson",           "CONV",            "–",                         "ECS 옵션 추가 예정"],
        ["ME",      "Santa Cruz",       "CONV / SELF",     "–",                         "ECS 옵션 추가 예정"],
        ["MVa",     "Tucson HEV",       "CONV / SELF",     "–",                         "ECS 옵션 추가 예정"],
        ["MEa",     "Santa Cruz HEV",   "SELF",            "Mando SFD3 (용량증대)",     "ECS 확대 · Tenneco 진입 기회"],
        ["MV GT",   "Tucson GT",        "ECS 100%",        "Mando 1 SOL SDC50",         "ECS 확대 · 북미 전용 금맥"],
    ]
    tbl = toss_table(s, Inches(0.6), Inches(1.95),
                     Inches(12.2), Inches(3.5),
                     data,
                     col_widths=[1.3, 2.2, 2.2, 3.0, 3.5],
                     header_fill=TEAL_DEEP,
                     header_color=WHITE,
                     body_size=11, header_size=10,
                     center_cols=[2])
    # Highlight MV GT row (ri=5) and MEa row (ri=4)
    for ri, fc in [(4, RGBColor(0xFF, 0xF8, 0xDC)), (5, TEAL_SOFT)]:
        for ci in range(5):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = fc
    # MV GT first col → TEAL
    mv_gt_cell = tbl.cell(5, 0)
    mv_gt_cell.fill.solid()
    mv_gt_cell.fill.fore_color.rgb = TEAL
    for p in mv_gt_cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = WHITE
    # MEa first col → LIME
    mea_cell = tbl.cell(4, 0)
    mea_cell.fill.solid()
    mea_cell.fill.fore_color.rgb = LIME

    # Bottom: Key strategic insight
    y_ins = Inches(5.7)
    # Hero callout — MV GT
    rounded(s, Inches(0.6), y_ins, Inches(7.5), Inches(1.3),
            fill=TEAL, radius=0.08)
    txt(s, Inches(0.85), y_ins + Inches(0.2),
        Inches(7), Inches(0.3),
        "TENNECO ATTACK POINT  ·  NORTH AMERICA ONLY",
        size=10, bold=True, color=LIME, font=FONT_EN)
    txt(s, Inches(0.85), y_ins + Inches(0.5),
        Inches(7), Inches(0.4),
        "MV GT = 북미 전용 100% ECS 라인업",
        size=16, bold=True, color=WHITE)
    txt(s, Inches(0.85), y_ins + Inches(0.92),
        Inches(7), Inches(0.35),
        "현재 Mando 독점 · PE 이후 ECS 확대 예정 → CVSA2로 전환 제안 최적 지점.",
        size=11, color=LIME)

    # Right: Supporting note
    rounded(s, Inches(8.3), y_ins, Inches(4.5), Inches(1.3),
            fill=TEAL_SOFT, radius=0.08)
    txt(s, Inches(8.55), y_ins + Inches(0.2),
        Inches(4.1), Inches(0.3),
        "MEa (Santa Cruz HEV)",
        size=11, bold=True, color=TEAL_DEEP)
    txt(s, Inches(8.55), y_ins + Inches(0.5),
        Inches(4.1), Inches(0.3),
        "Mando 용량증대 사양",
        size=14, bold=True, color=TEAL_DEEP)
    txt(s, Inches(8.55), y_ins + Inches(0.85),
        Inches(4.1), Inches(0.4),
        "PE 이후 ECS 확대 구간 — Tenneco 2nd-source 진입 가능성.",
        size=10, color=TEXT_BLACK, spacing=1.3)
    footer(s, 16)


# -------- Slide 17 — Q&A --------
def slide_17_qa():
    s = add_slide(WHITE)
    # Just a breath
    txt(s, 0, Inches(2.6), SW, Inches(2.0),
        "Q&A", size=160, bold=True, color=TEAL,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    txt(s, 0, Inches(5.0), SW, Inches(0.5),
        "궁금한 점을 나눠주세요.",
        size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    # tiny teal bar centered
    bar_w = Inches(0.6)
    rect(s, (SW - bar_w) / 2, Inches(5.8), bar_w, Inches(0.04),
         fill=TEAL)
    footer(s, 17)


# -------- Slide 18 — Thank You / Closing --------
def slide_18_closing():
    s = add_slide(WHITE)
    # tiny eyebrow top-left
    eyebrow(s, Inches(0.6), Inches(0.55),
            "TENNECO × HYUNDAI  ·  2026", color=TEAL, size=11)
    # Big headline
    txt(s, Inches(0.6), Inches(2.5), Inches(12), Inches(1.4),
        "Thank you.",
        size=100, bold=True, color=TEXT_BLACK, font=FONT_EN)
    # Thin bar
    rect(s, Inches(0.6), Inches(4.3), Inches(0.5), Inches(0.05),
         fill=TEAL)
    txt(s, Inches(0.6), Inches(4.45), Inches(12), Inches(0.35),
        "Tenneco × Hyundai Motor Group  ·  Next-Gen Chassis Partnership",
        size=15, color=TEXT_GRAY, font=FONT_EN)
    # Soft signature
    txt(s, Inches(0.6), Inches(4.95), Inches(12), Inches(0.3),
        "감사합니다.", size=16, color=TEAL_DEEP)

    # Bottom-right footer
    txt(s, Inches(9.5), Inches(7.1), Inches(3.3), Inches(0.3),
        "US HQ Internal Seminar  ·  2026.04",
        size=10, color=TEXT_LIGHT,
        align=PP_ALIGN.RIGHT, font=FONT_EN)
    txt(s, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25),
        "General Business – Tenneco Confidential",
        size=9, color=TEXT_LIGHT, italic=True)


# =======================================================================
#                               RUN
# =======================================================================
if __name__ == "__main__":
    slide_01_cover()
    slide_02_toc()
    slide_03_global_market()
    slide_04_america()
    slide_05_vs_competitors()
    slide_06_tenneco_ratio()
    slide_07_supply_overview()
    slide_08_volume_infographic()
    slide_09_market_status()
    slide_10_competitor_overview()
    slide_11_competitor_specificity()
    slide_12_competitor_matrix()
    slide_13_attack_strategy()
    slide_14_cvsa2()
    slide_15_ecs_lx3_lq2a()
    slide_16_ecs_mv_me()
    slide_17_qa()
    slide_18_closing()

    import os
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "Tenneco_Hyundai_Toss_KR.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")

