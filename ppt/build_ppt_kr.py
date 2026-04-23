# -*- coding: utf-8 -*-
"""
Tenneco × Hyundai Project Presentation (Korean version)
For Tenneco US HQ internal seminar - 16:9
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== Brand colors =====
TEAL = RGBColor(0x1A, 0x6F, 0x8C)
TEAL_DEEP = RGBColor(0x0E, 0x43, 0x57)
TEAL_DARK = RGBColor(0x12, 0x55, 0x6E)
LIME = RGBColor(0xD7, 0xE5, 0x5A)
ACCENT = RGBColor(0xE8, 0xA8, 0x3E)  # orange accent
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
BG_SOFT = RGBColor(0xE8, 0xF0, 0xF3)
TEXT_DARK = RGBColor(0x2B, 0x2B, 0x2B)
TEXT_MID = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GRAY_LINE = RGBColor(0xD0, 0xD0, 0xD0)
GRAY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)

FONT_KR = "맑은 고딕"
FONT_EN = "Arial"

# ===== Presentation setup =====
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ===== Helpers =====
def add_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w) if line_w else Pt(0.75)
    return shp


def add_round_rect(slide, x, y, w, h, fill=None, line=None, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_multi_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP):
    """runs = list of list of dicts (paragraphs of runs)
       each run dict: text, size, bold, color, font"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for rd in para:
            r = p.add_run()
            r.text = rd.get("text", "")
            r.font.name = rd.get("font", FONT_KR)
            r.font.size = Pt(rd.get("size", 14))
            r.font.bold = rd.get("bold", False)
            r.font.italic = rd.get("italic", False)
            c = rd.get("color", TEXT_DARK)
            r.font.color.rgb = c
    return tb


def add_line(slide, x1, y1, x2, y2, color=GRAY_LINE, weight=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_footer(slide, page_num=None, total=None):
    # Left footer
    add_text(slide, Inches(0.3), Inches(7.18), Inches(5), Inches(0.25),
             "General Business – Tenneco Confidential",
             size=9, color=TEXT_LIGHT, italic=True)
    # Right: Tenneco mark
    add_text(slide, Inches(11.2), Inches(7.18), Inches(2), Inches(0.25),
             "TENNECO × HYUNDAI", size=9, color=TEAL, bold=True,
             align=PP_ALIGN.RIGHT, font=FONT_EN)
    if page_num is not None:
        add_text(slide, Inches(12.7), Inches(7.18), Inches(0.55), Inches(0.25),
                 f"{page_num}", size=9, color=TEXT_MID,
                 align=PP_ALIGN.RIGHT, font=FONT_EN)


def add_top_bar(slide, section_num, section_name, title):
    """Thin top bar with section marker + title block for content slides."""
    # Top accent line
    add_rect(slide, 0, 0, SW, Inches(0.08), fill=TEAL_DEEP)
    # Section label tab
    add_rect(slide, Inches(0.5), Inches(0.22), Inches(1.1), Inches(0.38),
             fill=TEAL)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(1.1), Inches(0.38),
             f"SECTION {section_num:02d}", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # Section name
    add_text(slide, Inches(1.75), Inches(0.22), Inches(6), Inches(0.38),
             section_name, size=12, bold=False, color=TEXT_MID,
             anchor=MSO_ANCHOR.MIDDLE)
    # Main title
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
             title, size=26, bold=True, color=TEAL_DEEP,
             anchor=MSO_ANCHOR.MIDDLE)
    # Divider under title
    add_line(slide, Inches(0.5), Inches(1.45), Inches(12.83), Inches(1.45),
             color=LIME, weight=2)


def add_placeholder(slide, x, y, w, h, label="이미지 삽입 자리"):
    """Dashed-looking placeholder block for images to be added later."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_SOFT
    shp.line.color.rgb = TEAL
    shp.line.width = Pt(1.0)
    # dashed
    lnL = shp.line._get_or_add_ln()
    prstDash = etree.SubElement(lnL, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    # label text
    add_text(slide, x, y, w, h,
             f"📷  {label}",
             size=12, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shp


def add_bullet_block(slide, x, y, w, h, items, size=12, color=TEXT_DARK,
                     bullet="• ", spacing=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = bullet + it
        r.font.name = FONT_KR
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def make_table(slide, x, y, w, h, data, col_widths=None, header_fill=TEAL,
               header_color=WHITE, body_fill=WHITE, body_alt=BG_LIGHT,
               font_size=11, header_size=12, first_col_fill=None,
               first_col_color=None):
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            # fill
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif first_col_fill is not None and ci == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = first_col_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = body_alt if (ri % 2 == 0) else body_fill
            # text
            tf = cell.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ri == 0 else PP_ALIGN.LEFT
            p.text = ""
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT_KR
            if ri == 0:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = header_color
            else:
                r.font.size = Pt(font_size)
                if first_col_color is not None and ci == 0:
                    r.font.color.rgb = first_col_color
                    r.font.bold = True
                else:
                    r.font.color.rgb = TEXT_DARK
    return tbl


def add_kpi_card(slide, x, y, w, h, title, value, unit="", fill=TEAL,
                 title_color=WHITE, value_color=WHITE):
    add_round_rect(slide, x, y, w, h, fill=fill, radius=0.1)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.35), title,
             size=11, bold=False, color=title_color, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.48), w, Inches(0.6), value,
             size=28, bold=True, color=value_color, align=PP_ALIGN.CENTER,
             font=FONT_EN)
    if unit:
        add_text(slide, x, y + Inches(1.05), w, Inches(0.3), unit,
                 size=10, color=title_color, align=PP_ALIGN.CENTER)


def add_section_divider(slide, section_num, section_kr, section_en, page):
    # background full teal
    add_rect(slide, 0, 0, SW, SH, fill=TEAL_DEEP)
    # lime accent strip
    add_rect(slide, 0, Inches(3.4), SW, Inches(0.06), fill=LIME)
    # section num huge
    add_text(slide, Inches(0.8), Inches(2.1), Inches(4), Inches(1.8),
             f"{section_num:02d}", size=140, bold=True, color=LIME,
             font=FONT_EN, anchor=MSO_ANCHOR.MIDDLE)
    # Korean section name
    add_text(slide, Inches(4.8), Inches(2.6), Inches(8), Inches(0.8),
             section_kr, size=34, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.BOTTOM)
    # English subtitle
    add_text(slide, Inches(4.8), Inches(3.5), Inches(8), Inches(0.6),
             section_en, size=16, color=LIME, font=FONT_EN,
             anchor=MSO_ANCHOR.TOP)
    # Tenneco ribbon (bottom right)
    add_text(slide, Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4),
             "TENNECO × HYUNDAI", size=11, bold=True, color=LIME,
             align=PP_ALIGN.RIGHT, font=FONT_EN)
    # page
    add_text(slide, Inches(0.5), Inches(6.8), Inches(2), Inches(0.4),
             f"— {page} —", size=10, color=LIME, font=FONT_EN)


# ===== Flag / Map helpers =====
# Flag colors (RGB tuples) used as stacked colored bars to represent country flag
FLAGS = {
    "CN": [RGBColor(0xDE, 0x28, 0x10)],                                   # China — red (simplified)
    "IN": [RGBColor(0xFF, 0x99, 0x33), RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x13, 0x88, 0x08)],                                   # India (tricolor horizontal)
    "BR": [RGBColor(0x00, 0x9C, 0x3B), RGBColor(0xFE, 0xDF, 0x00)],       # Brazil (green/yellow)
    "KR": [RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0x00, 0x20, 0x71)],       # Korea (white/blue)
    "US": [RGBColor(0xB2, 0x22, 0x34), RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x3C, 0x3B, 0x6E)],                                   # USA
    "DE": [RGBColor(0x00, 0x00, 0x00), RGBColor(0xDD, 0x00, 0x00),
           RGBColor(0xFF, 0xCE, 0x00)],                                   # Germany
    "MX": [RGBColor(0x00, 0x69, 0x33), RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0xCE, 0x11, 0x26)],                                   # Mexico
    "AR": [RGBColor(0x74, 0xAC, 0xDF), RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0x74, 0xAC, 0xDF)],                                   # Argentina
    "CA": [RGBColor(0xFF, 0x00, 0x00), RGBColor(0xFF, 0xFF, 0xFF),
           RGBColor(0xFF, 0x00, 0x00)],                                   # Canada (simplified red-white-red)
}

def add_flag(slide, x, y, w, h, code):
    """Draw a simplified horizontal-stripe flag for the given country code."""
    colors = FLAGS.get(code)
    if not colors:
        return
    # outer border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    border.shadow.inherit = False
    border.fill.background()
    border.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    border.line.width = Pt(0.5)
    n = len(colors)
    stripe_h = Emu(int(h / n))
    for i, c in enumerate(colors):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x,
                                     y + Emu(int(h * i / n)), w, stripe_h)
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.fill.background()


def add_world_map_base(slide, x, y, w, h, fill=BG_SOFT, outline=TEAL):
    """Draw a very stylized world map base using rectangles for continents.
       This is a schematic, not geographically accurate."""
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()

    # rough continent blobs (normalized coordinates relative to x, y, w, h)
    # (cx_ratio, cy_ratio, w_ratio, h_ratio, name)
    continents = [
        # Americas
        (0.12, 0.32, 0.11, 0.22, "N.America"),
        (0.18, 0.68, 0.07, 0.28, "S.America"),
        # Europe
        (0.48, 0.26, 0.09, 0.16, "Europe"),
        # Africa
        (0.52, 0.55, 0.10, 0.30, "Africa"),
        # Asia
        (0.72, 0.32, 0.20, 0.30, "Asia"),
        # Oceania
        (0.86, 0.78, 0.10, 0.12, "Oceania"),
    ]
    for cxr, cyr, wr, hr, _ in continents:
        cx = x + Emu(int(w * cxr))
        cy = y + Emu(int(h * cyr))
        cw = Emu(int(w * wr))
        ch = Emu(int(h * hr))
        cont = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      cx - Emu(int(cw/2)),
                                      cy - Emu(int(ch/2)),
                                      cw, ch)
        cont.shadow.inherit = False
        cont.fill.solid()
        cont.fill.fore_color.rgb = RGBColor(0xC0, 0xD4, 0xDC)
        cont.line.color.rgb = outline
        cont.line.width = Pt(0.75)


def add_map_marker(slide, x, y, num=None, label=None, color=None,
                   size=0.32, text_color=None):
    """Draw a numbered/lettered circular marker on the map at (x,y).
       x, y are the CENTER of the marker."""
    if color is None:
        color = TEAL
    if text_color is None:
        text_color = WHITE
    r = Inches(size / 2)
    mk = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                x - r, y - r,
                                Inches(size), Inches(size))
    mk.shadow.inherit = False
    mk.fill.solid()
    mk.fill.fore_color.rgb = color
    mk.line.color.rgb = WHITE
    mk.line.width = Pt(1.5)
    if num is not None:
        add_text(slide, x - r, y - r, Inches(size), Inches(size),
                 str(num), size=10, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
    if label is not None:
        # small label to the right
        add_text(slide, x + r + Inches(0.05), y - Inches(0.13),
                 Inches(1.3), Inches(0.26),
                 label, size=9, bold=True, color=TEAL_DEEP, font=FONT_EN)


def add_na_map(slide, x, y, w, h, fill=BG_SOFT):
    """Draw a stylized North America outline (US + Canada + Mexico)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()

    # Canada (top wide)
    ca = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + Emu(int(w * 0.05)),
                                y + Emu(int(h * 0.05)),
                                Emu(int(w * 0.88)),
                                Emu(int(h * 0.30)))
    ca.shadow.inherit = False
    ca.adjustments[0] = 0.2
    ca.fill.solid()
    ca.fill.fore_color.rgb = RGBColor(0xC0, 0xD4, 0xDC)
    ca.line.color.rgb = TEAL
    ca.line.width = Pt(0.75)
    add_text(slide, x + Emu(int(w * 0.05)), y + Emu(int(h * 0.12)),
             Emu(int(w * 0.88)), Inches(0.3),
             "CANADA", size=10, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # USA (middle)
    us = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + Emu(int(w * 0.08)),
                                y + Emu(int(h * 0.38)),
                                Emu(int(w * 0.82)),
                                Emu(int(h * 0.32)))
    us.shadow.inherit = False
    us.adjustments[0] = 0.15
    us.fill.solid()
    us.fill.fore_color.rgb = RGBColor(0xA8, 0xC2, 0xCC)
    us.line.color.rgb = TEAL
    us.line.width = Pt(0.75)
    add_text(slide, x + Emu(int(w * 0.08)), y + Emu(int(h * 0.43)),
             Emu(int(w * 0.82)), Inches(0.3),
             "UNITED STATES", size=11, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # Mexico (bottom)
    mx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + Emu(int(w * 0.18)),
                                y + Emu(int(h * 0.73)),
                                Emu(int(w * 0.55)),
                                Emu(int(h * 0.22)))
    mx.shadow.inherit = False
    mx.adjustments[0] = 0.15
    mx.fill.solid()
    mx.fill.fore_color.rgb = RGBColor(0xC0, 0xD4, 0xDC)
    mx.line.color.rgb = TEAL
    mx.line.width = Pt(0.75)
    add_text(slide, x + Emu(int(w * 0.18)), y + Emu(int(h * 0.78)),
             Emu(int(w * 0.55)), Inches(0.3),
             "MEXICO", size=10, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)


def add_sa_map(slide, x, y, w, h, fill=BG_SOFT):
    """Stylized South America — Brazil + Argentina + surrounding."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()

    # Brazil (big, right side)
    br = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + Emu(int(w * 0.42)),
                                y + Emu(int(h * 0.12)),
                                Emu(int(w * 0.50)),
                                Emu(int(h * 0.55)))
    br.shadow.inherit = False
    br.adjustments[0] = 0.15
    br.fill.solid()
    br.fill.fore_color.rgb = RGBColor(0xA8, 0xC2, 0xCC)
    br.line.color.rgb = TEAL
    br.line.width = Pt(0.75)
    add_text(slide, x + Emu(int(w * 0.42)), y + Emu(int(h * 0.32)),
             Emu(int(w * 0.50)), Inches(0.3),
             "BRAZIL", size=12, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # Argentina (south, thin)
    ar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + Emu(int(w * 0.28)),
                                y + Emu(int(h * 0.65)),
                                Emu(int(w * 0.25)),
                                Emu(int(h * 0.30)))
    ar.shadow.inherit = False
    ar.adjustments[0] = 0.2
    ar.fill.solid()
    ar.fill.fore_color.rgb = RGBColor(0xC0, 0xD4, 0xDC)
    ar.line.color.rgb = TEAL
    ar.line.width = Pt(0.75)
    add_text(slide, x + Emu(int(w * 0.28)), y + Emu(int(h * 0.77)),
             Emu(int(w * 0.25)), Inches(0.3),
             "ARGENTINA", size=9, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # Other (Chile/Peru - left)
    other = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Emu(int(w * 0.08)),
                                   y + Emu(int(h * 0.18)),
                                   Emu(int(w * 0.20)),
                                   Emu(int(h * 0.55)))
    other.shadow.inherit = False
    other.adjustments[0] = 0.3
    other.fill.solid()
    other.fill.fore_color.rgb = RGBColor(0xC0, 0xD4, 0xDC)
    other.line.color.rgb = TEAL
    other.line.width = Pt(0.75)


# =======================================================================
#                           S L I D E S
# =======================================================================

# -------------------- S1: Cover --------------------
def slide_cover():
    s = add_slide()
    # background teal
    add_rect(s, 0, 0, SW, SH, fill=TEAL_DEEP)
    # diagonal accent band (faux)
    accent = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                Inches(-1), Inches(5.8), Inches(16), Inches(1.2))
    accent.shadow.inherit = False
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    accent.adjustments[0] = 0.6
    # Lime thin line
    add_rect(s, 0, Inches(5.65), SW, Inches(0.05), fill=LIME)
    # Top label
    add_text(s, Inches(0.8), Inches(0.8), Inches(12), Inches(0.4),
             "TENNECO INTERNAL SEMINAR — 2026", size=13, bold=True,
             color=LIME, font=FONT_EN)
    # Title Korean
    add_text(s, Inches(0.8), Inches(1.8), Inches(12), Inches(1.3),
             "Hyundai Motor Group 프로젝트 현황 및 공략 전략",
             size=40, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Title English
    add_text(s, Inches(0.8), Inches(3.1), Inches(12), Inches(0.7),
             "Tenneco × Hyundai : Current Business & Growth Strategy",
             size=20, color=LIME, font=FONT_EN, anchor=MSO_ANCHOR.MIDDLE)
    # Sub description
    add_text(s, Inches(0.8), Inches(4.0), Inches(12), Inches(0.5),
             "글로벌 시장 현황  ·  공급 프로그램  ·  경쟁사 비교  ·  CVSA2 확대  ·  운영 로드맵",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Presenter block (bottom)
    add_text(s, Inches(0.8), Inches(6.2), Inches(6), Inches(0.35),
             "Tenneco Korea – Hyundai/Kia Account", size=13, bold=True,
             color=WHITE, font=FONT_EN)
    add_text(s, Inches(0.8), Inches(6.55), Inches(6), Inches(0.3),
             "대한민국 · Tenneco 본사 세미나 발표 자료",
             size=11, color=LIME)
    # Right side big "T"
    add_text(s, Inches(10.8), Inches(1.2), Inches(2.2), Inches(2.5),
             "T", size=260, bold=True, color=TEAL, font=FONT_EN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# -------------------- S2: 목차 --------------------
def slide_toc():
    s = add_slide()
    add_rect(s, 0, 0, SW, SH, fill=WHITE)
    # Left teal column
    add_rect(s, 0, 0, Inches(4.3), SH, fill=TEAL_DEEP)
    add_text(s, Inches(0.5), Inches(0.6), Inches(3.6), Inches(0.4),
             "AGENDA", size=14, bold=True, color=LIME, font=FONT_EN)
    add_text(s, Inches(0.5), Inches(1.1), Inches(3.6), Inches(1.2),
             "목 차", size=66, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.5), Inches(2.7), Inches(3.6), Inches(0.6),
             "Hyundai Project\nPresentation", size=16, color=LIME,
             font=FONT_EN)
    add_text(s, Inches(0.5), Inches(6.7), Inches(3.6), Inches(0.5),
             "2026", size=22, bold=True, color=LIME, font=FONT_EN)

    items = [
        ("01", "현대 글로벌 시장 판매", "Hyundai Global Market Share"),
        ("02", "현대 북미 및 남미", "Hyundai N.America & S.America Market"),
        ("03", "현대에서의 Tenneco 점유율", "Tenneco Business Ratio in Hyundai"),
        ("04", "Tenneco의 현대 공급 프로그램", "Hyundai Supply Program (China · India · Brazil)"),
        ("05", "경쟁사 비교", "Competitor Analysis (Mando · ZF · Bilstein)"),
        ("06", "CVSA2 프로그램 확대 방안", "CVSA2 Program & Expansion Plan"),
        ("07", "현대 공략 방법", "Hyundai Winning Strategy"),
        ("08", "Tenneco Korea의 단계별 운영 방안", "Tenneco Korea Operating Roadmap"),
        ("09", "Q & A", "Discussion"),
    ]
    # Right content
    x0 = Inches(4.8)
    y0 = Inches(0.9)
    row_h = Inches(0.65)
    for i, (num, kr, en) in enumerate(items):
        y = y0 + row_h * i
        # number circle
        add_round_rect(s, x0, y + Inches(0.05), Inches(0.55), Inches(0.55),
                       fill=LIME, radius=0.5)
        add_text(s, x0, y + Inches(0.05), Inches(0.55), Inches(0.55),
                 num, size=16, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        # KR title
        add_text(s, x0 + Inches(0.8), y, Inches(7.5), Inches(0.38),
                 kr, size=17, bold=True, color=TEAL_DEEP)
        # EN subtitle
        add_text(s, x0 + Inches(0.8), y + Inches(0.38), Inches(7.5), Inches(0.28),
                 en, size=10, color=TEXT_MID, italic=True, font=FONT_EN)
        # divider line
        if i < len(items) - 1:
            add_line(s, x0 + Inches(0.8), y + row_h - Inches(0.03),
                     x0 + Inches(8.2), y + row_h - Inches(0.03),
                     color=GRAY_LINE, weight=0.5)
    add_footer(s, 2)


# -------------------- S3: Section 1 divider --------------------
def slide_section1():
    s = add_slide()
    add_section_divider(s, 1, "현대 글로벌 시장 판매",
                        "Hyundai Global Market Share · 2025", 3)


# -------------------- S4: Hyundai Global Market Share --------------------
def slide_global_market():
    s = add_slide()
    add_top_bar(s, 1, "현대 글로벌 시장 판매",
                "1. Hyundai Global Market Share  |  2025년 기준")
    # KPI cards row (top)
    add_kpi_card(s, Inches(0.5), Inches(1.7), Inches(2.8), Inches(1.3),
                 "2025 현대차 글로벌 판매", "414만", unit="대 (Total Sales)")
    add_kpi_card(s, Inches(3.45), Inches(1.7), Inches(2.8), Inches(1.3),
                 "최대 판매 지역", "아시아", unit="146만 대 (35.3%)",
                 fill=TEAL_DARK)
    add_kpi_card(s, Inches(6.4), Inches(1.7), Inches(2.8), Inches(1.3),
                 "핵심 전략 시장", "아메리카", unit="134만 대 (32.4%)",
                 fill=TEAL)
    add_kpi_card(s, Inches(9.35), Inches(1.7), Inches(2.8), Inches(1.3),
                 "기준 연도", "2025", unit="Year of Record",
                 fill=ACCENT, title_color=WHITE, value_color=WHITE)

    # Regional table (left)
    add_text(s, Inches(0.5), Inches(3.2), Inches(7.5), Inches(0.4),
             "■ 2025년 현대자동차 지역별 판매 현황", size=14, bold=True,
             color=TEAL_DEEP)
    table_data = [
        ["지역", "주요 국가 및 거점", "판매 수량 (추정)", "비중"],
        ["아메리카", "미국, 캐나다, 브라질(HMB), 멕시코", "약 134만 대", "32.4%"],
        ["아시아", "한국(내수), 인도, 베트남, 인도네시아", "약 146만 대", "35.3%"],
        ["유럽", "독일, 영국, 프랑스 등 EU 지역", "약 61만 대", "14.7%"],
        ["기타", "중동, 아프리카, 러시아, 호주 등", "약 73만 대", "17.6%"],
        ["합계", "전 세계 시장", "약 414만 대", "100%"],
    ]
    make_table(s, Inches(0.5), Inches(3.65), Inches(7.5), Inches(2.8),
               table_data, col_widths=[1.4, 3.5, 1.8, 1.0],
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # Right column: Tenneco footprint
    add_rect(s, Inches(8.3), Inches(3.2), Inches(4.55), Inches(3.3),
             fill=TEAL_DEEP)
    add_text(s, Inches(8.45), Inches(3.3), Inches(4.3), Inches(0.4),
             "TENNECO GLOBAL FOOTPRINT", size=12, bold=True,
             color=LIME, font=FONT_EN)
    add_text(s, Inches(8.45), Inches(3.65), Inches(4.3), Inches(0.35),
             "전 세계 공급 역량", size=11, color=WHITE)

    kpis = [
        ("60,000", "Global Team Members"),
        ("184", "Manufacturing Plants"),
        ("23", "Distribution Centers"),
        ("40", "Engineering Facilities"),
    ]
    for i, (v, lbl) in enumerate(kpis):
        yy = Inches(4.1 + i * 0.55)
        add_text(s, Inches(8.45), yy, Inches(1.3), Inches(0.4),
                 v, size=22, bold=True, color=LIME,
                 font=FONT_EN, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(9.85), yy + Inches(0.08), Inches(2.9),
                 Inches(0.4), lbl, size=11, color=WHITE, font=FONT_EN)

    # Bottom insight
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.5),
             fill=BG_SOFT)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.5),
             "Insight  ·  현대차는 지정학적 리스크와 경기 변동 속에서도 북미 최대 실적 및 인도 시장 성장을 통해 글로벌 판매 비중을 성공적으로 다변화했다.",
             size=11, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE, italic=True)
    add_footer(s, 4)


# -------------------- S5: Section 2 --------------------
def slide_section2():
    s = add_slide()
    add_section_divider(s, 2, "현대 북미 및 남미",
                        "Hyundai N.America & S.America Market · 2025", 5)


# -------------------- S6: 2-1 America detail --------------------
def slide_america_detail():
    s = add_slide()
    add_top_bar(s, 2, "현대 북미 및 남미",
                "2-1. Hyundai N.America & S.America Market  |  국가별 판매 현황")

    # Left: table
    add_text(s, Inches(0.5), Inches(1.7), Inches(7), Inches(0.4),
             "■ 2025년 현대자동차 핵심 국가별 판매 현황 (글로벌 134만 대 기준)",
             size=13, bold=True, color=TEAL_DEEP)
    data = [
        ["지역", "핵심 국가", "판매 수량 (추정)", "글로벌 비중"],
        ["북미", "미국 (USA)", "약 90.2만 대", "21.8%"],
        ["북미", "캐나다 (Canada)", "약 14.1만 대", "3.4%"],
        ["북미", "멕시코 (Mexico)", "약 5.1만 대", "1.2%"],
        ["남미", "브라질 (HMB)", "약 18.5만 대", "4.5%"],
        ["기타", "기타 아메리카 국가", "약 6.1만 대", "1.5%"],
        ["합계", "아메리카 전체", "약 134.0만 대", "32.4%"],
    ]
    make_table(s, Inches(0.5), Inches(2.15), Inches(7.2), Inches(3.7), data,
               col_widths=[1.1, 2.9, 2.0, 1.2],
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP,
               font_size=11)

    # Right: bar chart rendered with rectangles
    add_text(s, Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
             "■ 판매량 시각화 (만 대)", size=13, bold=True, color=TEAL_DEEP)
    bars = [
        ("미국", 90.2, TEAL_DEEP),
        ("브라질", 18.5, TEAL),
        ("캐나다", 14.1, TEAL),
        ("기타", 6.1, ACCENT),
        ("멕시코", 5.1, ACCENT),
    ]
    bar_x = Inches(9.2)
    bar_max_w = Inches(3.4)
    bar_max_val = 100.0
    by = Inches(2.2)
    for i, (lbl, val, col) in enumerate(bars):
        yy = by + Inches(0.6) * i
        add_text(s, Inches(8.0), yy, Inches(1.15), Inches(0.4),
                 lbl, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.RIGHT)
        bw = Emu(int(bar_max_w * val / bar_max_val))
        add_rect(s, bar_x, yy + Inches(0.08), bw, Inches(0.28), fill=col)
        add_text(s, bar_x + bw + Inches(0.05), yy, Inches(1), Inches(0.4),
                 f"{val:.1f}만", size=10, color=TEXT_MID,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

    # Bottom insight
    add_rect(s, Inches(0.5), Inches(6.15), Inches(12.33), Inches(0.9),
             fill=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.35),
             "KEY TAKEAWAY", size=10, bold=True, color=LIME, font=FONT_EN)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5),
             "아메리카 판매의 67%가 미국에 집중 (90.2만 대). 브라질·멕시코의 현지 생산 기반(HMB)은 Tenneco 공급 확장의 핵심 레버리지.",
             size=12, color=WHITE, anchor=MSO_ANCHOR.TOP)
    add_footer(s, 6)


# -------------------- S7: 2-2 Competitor comparison --------------------
def slide_america_vs_competitor():
    s = add_slide()
    add_top_bar(s, 2, "현대 북미 및 남미",
                "2-2. GM · Toyota · Hyundai 판매량 비교  |  2025년 기준")

    # comparison cards (3 big cards)
    data = [
        {"name": "GM", "vol": "약 320 ~ 340만 대", "rank": "북미 1위",
         "desc": "대형 픽업트럭과 SUV 중심의 압도적 시장 지배력",
         "fill": TEAL_DEEP, "tag": "DOMINANT LEADER"},
        {"name": "TOYOTA", "vol": "약 280 ~ 300만 대", "rank": "하이브리드 1위",
         "desc": "하이브리드 절대 강자. 북미·남미 모두에서 높은 신뢰도 확보",
         "fill": TEAL, "tag": "HYBRID POWERHOUSE"},
        {"name": "HYUNDAI", "vol": "약 134만 대", "rank": "성장세 1위",
         "desc": "전기차 및 SUV 라인업 강화를 통해 추격 중. Tenneco의 주력 고객",
         "fill": ACCENT, "tag": "FASTEST GROWING"},
    ]
    card_w = Inches(4.1)
    gap = Inches(0.1)
    x0 = Inches(0.5)
    y0 = Inches(1.8)
    card_h = Inches(3.8)
    for i, d in enumerate(data):
        x = x0 + (card_w + gap) * i
        add_rect(s, x, y0, card_w, card_h, fill=d["fill"])
        # Tag
        add_text(s, x + Inches(0.2), y0 + Inches(0.2), card_w - Inches(0.4),
                 Inches(0.3), d["tag"], size=10, bold=True, color=LIME,
                 font=FONT_EN)
        # Name
        add_text(s, x + Inches(0.2), y0 + Inches(0.55), card_w - Inches(0.4),
                 Inches(0.8), d["name"], size=36, bold=True, color=WHITE,
                 font=FONT_EN)
        # Line
        add_rect(s, x + Inches(0.2), y0 + Inches(1.45),
                 Inches(0.8), Inches(0.04), fill=LIME)
        # Volume label
        add_text(s, x + Inches(0.2), y0 + Inches(1.6), card_w - Inches(0.4),
                 Inches(0.3), "아메리카 총 판매량", size=10, color=LIME)
        add_text(s, x + Inches(0.2), y0 + Inches(1.9), card_w - Inches(0.4),
                 Inches(0.45), d["vol"], size=20, bold=True, color=WHITE)
        # Rank
        add_text(s, x + Inches(0.2), y0 + Inches(2.45), card_w - Inches(0.4),
                 Inches(0.3), "시장 지위", size=10, color=LIME)
        add_text(s, x + Inches(0.2), y0 + Inches(2.7), card_w - Inches(0.4),
                 Inches(0.35), d["rank"], size=15, bold=True, color=WHITE)
        # Desc
        add_text(s, x + Inches(0.2), y0 + Inches(3.1), card_w - Inches(0.4),
                 Inches(0.7), d["desc"], size=10, color=WHITE)
    # bottom strip
    add_rect(s, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2),
             fill=BG_SOFT)
    add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.35),
             "■ 전략적 해석", size=12, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.7),
             "현대차는 GM·Toyota 대비 절대 판매량은 낮지만 성장률·전동화 속도·신규 플랫폼(EV GMP) 기준으로 가장 공격적. "
             "Tenneco 입장에서 중장기 볼륨 확대 여지가 가장 큰 OEM.",
             size=11, color=TEXT_DARK)
    add_footer(s, 7)


# -------------------- S8: Section 3 --------------------
def slide_section3():
    s = add_slide()
    add_section_divider(s, 3, "현대에서의 Tenneco 점유율",
                        "Tenneco Business Ratio in Hyundai", 8)


# -------------------- S9: Tenneco share in Hyundai --------------------
def slide_tenneco_share():
    s = add_slide()
    add_top_bar(s, 3, "현대에서의 Tenneco 점유율",
                "3. Tenneco Business Ratio in Hyundai  |  공급사별 점유 현황")

    # Left: pie-chart-like visualization using rectangles (stacked)
    add_text(s, Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
             "■ 현대/기아 섀시 댐퍼 공급사 점유율 (추정)", size=13, bold=True,
             color=TEAL_DEEP)
    # Horizontal stacked bar
    bar_x = Inches(0.5)
    bar_y = Inches(2.3)
    bar_w = Inches(6.2)
    bar_h = Inches(0.6)
    segs = [
        ("Mando", 65, TEAL_DEEP),
        ("ZF", 18, TEAL),
        ("Bilstein", 10, ACCENT),
        ("Tenneco", 7, LIME),
    ]
    cur_x = bar_x
    for name, pct, col in segs:
        w = Emu(int(bar_w * pct / 100))
        add_rect(s, cur_x, bar_y, w, bar_h, fill=col)
        tc = TEAL_DEEP if col == LIME else WHITE
        add_text(s, cur_x, bar_y, w, bar_h,
                 f"{name}\n{pct}%", size=11, bold=True, color=tc,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x = Emu(cur_x + w)
    # legend note
    add_text(s, Inches(0.5), Inches(3.05), Inches(6.2), Inches(0.35),
             "* 추정치: Mando 국내 현대/기아 기본 섀시 독점 + ZF(Genesis/Ioniq6 N) · Bilstein(High-end) · Tenneco(초기 진입)",
             size=9, italic=True, color=TEXT_LIGHT)

    # Supplier breakdown list under bar
    break_data = [
        ("Mando", "현대/기아 기본 섀시의 대부분 (ICE/EV 공용 Conventional + Semi-active)"),
        ("ZF", "Genesis Entry-Level, Ioniq6 N (Semi-active ECS)"),
        ("Bilstein", "Genesis High-end (Air shock + ECS) 전용 포지션"),
        ("Tenneco", "초기 진입 단계 — CVSA2 Demo Car 검증 단계"),
    ]
    add_text(s, Inches(0.5), Inches(3.5), Inches(6.2), Inches(0.35),
             "■ 공급 영역 Summary", size=13, bold=True, color=TEAL_DEEP)
    yy = Inches(3.95)
    for name, desc in break_data:
        col = {"Mando": TEAL_DEEP, "ZF": TEAL, "Bilstein": ACCENT,
               "Tenneco": LIME}[name]
        add_rect(s, Inches(0.5), yy, Inches(0.15), Inches(0.4), fill=col)
        add_text(s, Inches(0.75), yy, Inches(1.4), Inches(0.4),
                 name, size=11, bold=True, color=TEAL_DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        add_text(s, Inches(2.15), yy, Inches(4.7), Inches(0.4),
                 desc, size=10, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        yy = Emu(yy + Inches(0.45))

    # Right: Opportunity call-out
    add_rect(s, Inches(7.0), Inches(1.7), Inches(5.83), Inches(4.8),
             fill=TEAL_DEEP)
    add_text(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.35),
             "OPPORTUNITY", size=11, bold=True, color=LIME, font=FONT_EN)
    add_text(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.7),
             "Tenneco의 점유 확장 여지",
             size=22, bold=True, color=WHITE)
    add_line(s, Inches(7.2), Inches(3.0), Inches(9.0), Inches(3.0),
             color=LIME, weight=2)

    ops = [
        ("▲ 2-Valve MacPherson", "현대가 차기 제품에 요구하는 핵심 기술 — Tenneco만 보유"),
        ("▲ 가격 경쟁력", "ZF·Bilstein 대비 HMC가 합리적이라 평가"),
        ("▲ Golf 유럽 시승 호평", "2025년 4월 HMC 엔지니어 긍정 반응"),
        ("▲ VP Manfred Harrer", "유럽 공급사 선호 기조 확립"),
        ("▲ N → Genesis 수직 확산", "N Brand 검증 시 Genesis 자동 이관 구조"),
    ]
    yy = Inches(3.2)
    for t, d in ops:
        add_text(s, Inches(7.2), yy, Inches(5.5), Inches(0.3),
                 t, size=13, bold=True, color=LIME)
        add_text(s, Inches(7.35), yy + Inches(0.3), Inches(5.3), Inches(0.35),
                 d, size=10, color=WHITE)
        yy = Emu(yy + Inches(0.67))

    # bottom call
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.45),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
             "현재 Tenneco 점유율은 한 자릿수 초기 진입. CVSA2 Ioniq6 N RFQ 확정 시 두 자릿수 급상승 가능.",
             size=12, bold=True, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 9)


# -------------------- S10: Section 4 divider --------------------
def slide_section4():
    s = add_slide()
    add_section_divider(s, 4, "테네코에서의 현대 공급 프로그램",
                        "Tenneco Hyundai Supply Programs  |  China · India · Brazil", 10)


# -------------------- S11: 4-1 Regional Overview --------------------
def slide_supply_overview():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-1.  지역별 프로젝트 카드  |  Country × Vehicle × Project")

    # Project cards: (country code, flag name, project code, vehicle name, year, marker_num)
    projects = [
        ("CN", "CHINA",  "GC",       "IX25",           "2014", 1),
        ("CN", "CHINA",  "YC FL",    "Verna",          "2019", 2),
        ("CN", "CHINA",  "ADC",      "Elantra 6th",    "2016", 3),
        ("CN", "CHINA",  "CN7c",     "Elantra 7th",    "2020", 4),
        ("IN", "INDIA",  "QXI",      "Venue",          "2019", 5),
        ("IN", "INDIA",  "BI3",      "Elite i20",      "2020", 6),
        ("BR", "BRAZIL", "BR2",      "HB20",           "2022", 7),
        ("BR", "BRAZIL", "SU2b",     "Creta",          "2023", 8),
        ("BR", "BRAZIL", "BC4b",     "(개발중)",        "2026", 9),
    ]
    # 3 rows x 3 cols of cards  (small, uniform)
    card_w = Inches(1.95)
    card_h = Inches(1.55)
    gap_x = Inches(0.12)
    gap_y = Inches(0.1)
    start_x = Inches(0.5)
    start_y = Inches(1.75)

    for i, (code, country, proj, veh, yr, num) in enumerate(projects):
        col = i % 3
        row = i // 3
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        # card body
        add_rect(s, x, y, card_w, card_h, fill=BG_LIGHT, line=GRAY_LINE, line_w=0.5)
        # flag + country (left top)
        add_flag(s, x + Inches(0.12), y + Inches(0.12), Inches(0.32),
                 Inches(0.2), code)
        add_text(s, x + Inches(0.5), y + Inches(0.1), Inches(1),
                 Inches(0.2),
                 country, size=8, bold=True, color=TEAL_DEEP, font=FONT_EN)
        # marker num (right top)
        add_map_marker(s, x + card_w - Inches(0.25),
                       y + Inches(0.22), num=num, color=TEAL, size=0.28)
        # vehicle placeholder
        add_rect(s, x + Inches(0.12), y + Inches(0.4),
                 card_w - Inches(0.24), Inches(0.55),
                 fill=BG_SOFT, line=TEAL, line_w=0.5)
        add_text(s, x + Inches(0.12), y + Inches(0.4),
                 card_w - Inches(0.24), Inches(0.55),
                 "🚗  차량 사진", size=9, color=TEAL,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # project code (big)
        add_text(s, x + Inches(0.12), y + Inches(1.0),
                 card_w - Inches(0.24), Inches(0.3),
                 proj, size=13, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.LEFT, font=FONT_EN)
        # vehicle & year (small)
        add_text(s, x + Inches(0.12), y + Inches(1.28),
                 card_w - Inches(0.24), Inches(0.2),
                 f"{veh}  ·  SOP {yr}", size=8, color=TEXT_MID,
                 align=PP_ALIGN.LEFT)

    # Right side : mini world map with markers
    map_x = Inches(6.95)
    map_y = Inches(1.75)
    map_w = Inches(5.85)
    map_h = Inches(4.75)
    add_world_map_base(s, map_x, map_y, map_w, map_h)
    # Title on map
    add_text(s, map_x, map_y + Inches(0.05), map_w, Inches(0.3),
             "GLOBAL PROJECT MAP",
             size=10, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    # Highlight markers (numbers at approximate country positions)
    # normalized positions within map
    positions = {
        "CN": (0.78, 0.38),   # China - Asia east
        "IN": (0.70, 0.48),   # India - South Asia
        "BR": (0.23, 0.70),   # Brazil - S.America
    }
    # Numbers grouped per country
    groups = {"CN": [1, 2, 3, 4], "IN": [5, 6], "BR": [7, 8, 9]}
    colors_per = {"CN": TEAL, "IN": TEAL_DARK, "BR": TEAL_DEEP}
    for code, (rx, ry) in positions.items():
        cx = map_x + Emu(int(map_w * rx))
        cy = map_y + Emu(int(map_h * ry))
        # country ring
        big = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - Inches(0.45), cy - Inches(0.45),
                                 Inches(0.9), Inches(0.9))
        big.shadow.inherit = False
        big.fill.solid()
        big.fill.fore_color.rgb = LIME
        big.line.color.rgb = colors_per[code]
        big.line.width = Pt(2)
        # country label
        add_text(s, cx - Inches(0.6), cy + Inches(0.5),
                 Inches(1.2), Inches(0.28),
                 code, size=11, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # number group
        nums_txt = "·".join(str(n) for n in groups[code])
        add_text(s, cx - Inches(0.45), cy - Inches(0.2),
                 Inches(0.9), Inches(0.4),
                 nums_txt, size=12, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)

    # Legend bar
    add_rect(s, map_x, map_y + map_h - Inches(0.6), map_w, Inches(0.5),
             fill=TEAL_DEEP)
    add_text(s, map_x + Inches(0.15), map_y + map_h - Inches(0.55),
             map_w - Inches(0.3), Inches(0.2),
             "PROJECT LOCATIONS",
             size=9, bold=True, color=LIME, font=FONT_EN)
    add_text(s, map_x + Inches(0.15), map_y + map_h - Inches(0.34),
             map_w - Inches(0.3), Inches(0.2),
             "중국 4 · 인도 2 · 브라질 3   (총 9개 프로젝트 / 3대륙)",
             size=10, color=WHITE)

    # bottom summary
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.45),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
             "총 9개 현대 프로그램 · 3대륙 동시 대응 — Tenneco 현대 공급 전략의 지리적 기반",
             size=12, bold=True, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 11)


# -------------------- S12: 4-2 China Project List --------------------
def slide_china_project():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-2.  China Project List  |  중국 프로젝트 5종")

    data = [
        ["Project Code", "Vehicle", "SOP", "EOP", "비고"],
        ["GC",           "IX25",             "2014.10", "2017.10", "1세대 CUV"],
        ["GC PE",        "IX25 (PE)",        "2017.10", "2019.05", "Phase Enhancement"],
        ["YC",           "Verna",            "2016.10", "2019.10", "B세그먼트 세단"],
        ["YC FL",        "Verna (FL)", "2019.10", "2022.03", "Face-lift"],
        ["ADC",          "Elantra 领动",     "2016.03", "2019.03", "6세대 세단"],
        ["ADC PE",       "Elantra 领动 PE",  "2019.03", "2020.10", "Phase Enhancement"],
        ["ADC PHEV",     "Elantra Hybrid",   "2019.05", "–",       "하이브리드"],
        ["CN7c",         "Elantra 7th",      "2020.10", "2023",    "7세대 현행"],
    ]
    make_table(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.7),
               data, col_widths=[1.4, 2.4, 1.5, 1.5, 3.0],
               font_size=12, header_size=12,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # note
    add_rect(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
             fill=BG_SOFT)
    add_text(s, Inches(0.8), Inches(6.55), Inches(12), Inches(0.45),
             "중국은 Tenneco의 현대 공급 베이스가 가장 두터운 지역. 2014년부터 연속 세대 공급 이어옴.",
             size=11, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 12)


# -------------------- S13: 4-3 India & Brazil Project List --------------------
def slide_india_brazil_project():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-3.  India & Brazil Project List")

    # India section
    add_rect(s, Inches(0.6), Inches(1.75), Inches(5.95), Inches(0.45),
             fill=TEAL)
    add_text(s, Inches(0.8), Inches(1.75), Inches(5.95), Inches(0.45),
             "INDIA  |  인도 — 2개 프로젝트",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_EN)
    india_data = [
        ["Code", "Vehicle", "SOP", "비고"],
        ["QXI",  "Venue",   "2019.05", "소형 SUV"],
        ["BI3",  "Elite i20", "2020.10", "해치백"],
    ]
    make_table(s, Inches(0.6), Inches(2.25), Inches(5.95), Inches(2.0),
               india_data, col_widths=[1.2, 2.5, 1.3, 2.0],
               font_size=12, first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # Brazil section
    add_rect(s, Inches(6.75), Inches(1.75), Inches(5.95), Inches(0.45),
             fill=TEAL_DEEP)
    add_text(s, Inches(6.95), Inches(1.75), Inches(5.95), Inches(0.45),
             "BRAZIL  |  브라질 — 3개 프로젝트",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_EN)
    brazil_data = [
        ["Code",    "Vehicle", "SOP",   "비고"],
        ["BR2",     "HB20",    "2022.06", "해치백 (현지모델)"],
        ["SU2b",    "Creta",   "2023~",  "소형 SUV"],
        ["SU2b PE", "Creta PE", "2024.12", "Phase Enhancement"],
        ["BC4b",    "–",       "2026~",  "개발중 (In development)"],
    ]
    make_table(s, Inches(6.75), Inches(2.25), Inches(5.95), Inches(3.0),
               brazil_data, col_widths=[1.2, 2.2, 1.3, 2.3],
               font_size=11, first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # Key insight
    add_rect(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.2),
             fill=BG_SOFT)
    add_rect(s, Inches(0.6), Inches(5.8), Inches(0.12), Inches(1.2),
             fill=LIME)
    add_text(s, Inches(0.9), Inches(5.9), Inches(11.8), Inches(0.4),
             "💡  Key Insight",
             size=12, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.9), Inches(6.25), Inches(11.8), Inches(0.8),
             "신흥시장 (인도·브라질) 중심으로 현대 현지 생산 모델 지속 수주 중.\n" +
             "BC4b 2026년 양산 예정 → 브라질 공급량 추가 확대 전망.",
             size=11, color=TEXT_DARK)
    add_footer(s, 13)


# -------------------- S14: 4-4 Program Vehicle Volume --------------------
def slide_program_volume():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-4.  Program Vehicle Volume  |  Sales by country infographic")

    # === LEFT: Compact table ===
    data = [
        ["지역", "Project", "2023", "2024", "2025", "2026", "2027"],
        ["🇮🇳 India",  "BI3",   "81,600",  "80,450",  "80,710",  "6,860",  "–"],
        ["🇮🇳 India",  "Qxi",   "138,300", "116,370", "57,850",  "–",      "–"],
        ["🇧🇷 Brazil", "BR2",   "137,600", "148,800", "145,800", "30,500", "–"],
        ["🇧🇷 Brazil", "SU2b",  "–",       "–",       "70,200",  "73,000", "43,000"],
        ["🇧🇷 Brazil", "BC4b",  "–",       "–",       "–",       "112,500","138,000"],
        ["합계(추정)", "",      "357,500", "345,620", "354,560", "222,860","181,000"],
    ]
    tbl = make_table(s, Inches(0.4), Inches(1.75), Inches(7.3), Inches(3.6),
                     data, col_widths=[1.1, 1.0, 0.84, 0.84, 0.84, 0.84, 0.84],
                     font_size=9, header_size=10,
                     first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)
    # Highlight totals row (last)
    try:
        last = tbl.rows[6]
        for c in last.cells:
            c.fill.solid()
            c.fill.fore_color.rgb = TEAL
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = WHITE
                    r.font.bold = True
    except Exception:
        pass

    # === RIGHT: Country ranking infographic with map markers ===
    # Title ribbon
    add_rect(s, Inches(7.9), Inches(1.75), Inches(4.9), Inches(0.38),
             fill=TEAL_DEEP)
    add_text(s, Inches(7.9), Inches(1.75), Inches(4.9), Inches(0.38),
             "Volume by Country (2023~2027 누적)",
             size=10, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Country ranks (simplified totals)
    # Brazil total = 137600+148800+145800+30500 + 70200+73000+43000 + 112500+138000 = 899,400
    # India total = 81600+80450+80710+6860 + 138300+116370+57850 = 562,140
    ranks = [
        ("01", "🇧🇷", "Brazil", "899,400",  "약 90만 대",  "BR2 · SU2b · BC4b", TEAL_DEEP),
        ("02", "🇮🇳", "India",  "562,140",  "약 56만 대",  "BI3 · Qxi",         TEAL),
    ]
    y0 = Inches(2.2)
    for i, (rk, flag, name, units, units_kr, projs, color) in enumerate(ranks):
        yy = y0 + Inches(1.55 * i)
        # row card
        add_rect(s, Inches(7.9), yy, Inches(4.9), Inches(1.45),
                 fill=BG_LIGHT)
        # rank tag
        add_rect(s, Inches(7.9), yy, Inches(0.7), Inches(1.45), fill=color)
        add_text(s, Inches(7.9), yy, Inches(0.7), Inches(0.7),
                 rk, size=22, bold=True, color=LIME,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_text(s, Inches(7.9), yy + Inches(0.75), Inches(0.7),
                 Inches(0.65),
                 flag, size=22, bold=False, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # country name
        add_text(s, Inches(8.7), yy + Inches(0.1), Inches(2.5),
                 Inches(0.35),
                 name, size=15, bold=True, color=TEAL_DEEP, font=FONT_EN)
        # units big number
        add_text(s, Inches(8.7), yy + Inches(0.45), Inches(2.5),
                 Inches(0.5),
                 units, size=22, bold=True, color=color, font=FONT_EN)
        add_text(s, Inches(8.7), yy + Inches(0.95), Inches(2.5),
                 Inches(0.3),
                 units_kr + " (5년 누적)", size=9, color=TEXT_MID)
        # projects
        add_rect(s, Inches(11.3), yy + Inches(0.25), Inches(1.4),
                 Inches(0.4), fill=color)
        add_text(s, Inches(11.3), yy + Inches(0.25), Inches(1.4),
                 Inches(0.4),
                 "Projects", size=9, bold=True, color=LIME,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_text(s, Inches(11.3), yy + Inches(0.7), Inches(1.4),
                 Inches(0.6),
                 projs, size=9, color=TEAL_DEEP, align=PP_ALIGN.CENTER,
                 bold=True)

    # === BOTTOM: World map with markers (2 countries) ===
    map_x = Inches(0.4)
    map_y = Inches(5.5)
    map_w = Inches(12.5)
    map_h = Inches(1.55)
    add_world_map_base(s, map_x, map_y, map_w, map_h)

    # Brazil marker (#01)
    br_x = map_x + Emu(int(map_w * 0.23))
    br_y = map_y + Emu(int(map_h * 0.70))
    br_big = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                br_x - Inches(0.25), br_y - Inches(0.25),
                                Inches(0.5), Inches(0.5))
    br_big.shadow.inherit = False
    br_big.fill.solid()
    br_big.fill.fore_color.rgb = TEAL_DEEP
    br_big.line.color.rgb = LIME
    br_big.line.width = Pt(2)
    add_text(s, br_x - Inches(0.25), br_y - Inches(0.25), Inches(0.5),
             Inches(0.5),
             "01", size=12, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_text(s, br_x - Inches(0.85), br_y + Inches(0.3),
             Inches(1.7), Inches(0.25),
             "BRAZIL · 899K", size=9, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # India marker (#02)
    in_x = map_x + Emu(int(map_w * 0.70))
    in_y = map_y + Emu(int(map_h * 0.48))
    in_big = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                in_x - Inches(0.25), in_y - Inches(0.25),
                                Inches(0.5), Inches(0.5))
    in_big.shadow.inherit = False
    in_big.fill.solid()
    in_big.fill.fore_color.rgb = TEAL
    in_big.line.color.rgb = LIME
    in_big.line.width = Pt(2)
    add_text(s, in_x - Inches(0.25), in_y - Inches(0.25), Inches(0.5),
             Inches(0.5),
             "02", size=12, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_text(s, in_x - Inches(0.8), in_y + Inches(0.3),
             Inches(1.6), Inches(0.25),
             "INDIA · 562K", size=9, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # Key takeaway ribbon at very bottom
    add_rect(s, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.08),
             fill=LIME)
    add_footer(s, 14)


# -------------------- S15: 4-5 China Part Spec 1 --------------------
def slide_partspec_china1():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-5.  China 부품 사양 ①  |  GC · YC FL · ADc")

    data = [
        ["Project", "F/R", "Rod OD", "Inner Tube", "Outer Tube", "Valve", "Part No."],
        ["GC",      "Front", "Ø22",       "Ø32×1.03T", "Ø51×2.5T",  "MTV",    "MEM184501/601"],
        ["GC",      "Rear",  "Ø12.4",     "Ø30×1.03T", "Ø45.3×1.27T","MTV",   "MEM184101/201"],
        ["YC FL",   "Front", "Ø20",       "Ø30×1.03T", "Ø48.6×2.5T","MTV",    "MET1W1101/1001"],
        ["YC FL",   "Rear",  "Ø12.4",     "Ø25.4×1.18T","Ø38.5×1.27T","MTV",  "MET1W0901"],
        ["ADc",     "Front", "Ø22 (hollow)","Ø32×1.03T","Ø50×2.5T",  "MTV CL", "MEM261601/701"],
        ["ADc",     "Rear",  "Ø12.4",     "Ø25.4×1.18T","Ø38.5×1.27T","MTV CL","MEM261901"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.2),
               data, col_widths=[1.3, 0.9, 1.5, 2.0, 2.3, 1.3, 3.0],
               font_size=11, header_size=11,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # legend
    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
             fill=BG_SOFT)
    add_text(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.35),
             "MTV : Mechanically Tuned Valve       MTV CL : MTV Closed-loop       OD : Outer Diameter       T : Thickness",
             size=10, bold=True, color=TEAL_DEEP, font=FONT_EN)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.35),
             "주요 사양: Front Ø22 / Rear Ø12.4 Rod · 2.5T 외경 튜브 · MTV 계열 밸브 (중국 표준)",
             size=10, color=TEXT_MID, italic=True)
    add_footer(s, 15)


# -------------------- S16: 4-6 China Part Spec 2 --------------------
def slide_partspec_china2():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-6.  China 부품 사양 ②  |  CN7c · ADc PHEV · NU2")

    data = [
        ["Project", "F/R", "Rod OD", "Inner Tube", "Outer Tube", "Valve", "Part No."],
        ["CN7c",      "Front", "Ø22 (hollow)", "Ø32×1.03T", "Ø50×2.5T",   "MTV CL", "MEM548001/549001"],
        ["CN7c",      "Rear",  "Ø12.4",        "Ø25.4×1.18T","Ø38.5×1.27T","MTV CL", "MEM585801"],
        ["ADc PHEV",  "Front", "Ø22 (hollow)", "Ø32×1.03T", "Ø50×2.5T",   "MTV CL", "MET1P7901/8001"],
        ["ADc PHEV",  "Rear",  "Ø12.4",        "Ø30×1.03T", "Ø45.3×1.27T","MTV CL", "MET1P8101"],
        ["NU2",       "Front", "Ø22",          "Ø35×1.03T", "Ø52×2.5T",   "MTV CL", "MX5197F26/F31"],
        ["NU2",       "Rear",  "Ø12.4",        "Ø30×1.03T", "Ø45.3×1.27T","MTV CL", "MX5198F07"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.0),
               data, col_widths=[1.3, 0.9, 1.6, 2.0, 2.3, 1.3, 2.9],
               font_size=11, header_size=11,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # highlight
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
             fill=BG_SOFT)
    add_rect(s, Inches(0.5), Inches(6.0), Inches(0.12), Inches(1.0),
             fill=LIME)
    add_text(s, Inches(0.8), Inches(6.1), Inches(12), Inches(0.35),
             "📌  7세대 Elantra(CN7c) 이후 프로젝트 공통 특징",
             size=12, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.8), Inches(6.45), Inches(12), Inches(0.5),
             "전량 MTV CL(Closed-Loop) 밸브로 전환 · Front Ø22 hollow rod 표준화 · Rear Ø45.3×1.27T 공용 플랫폼",
             size=11, color=TEXT_DARK)
    add_footer(s, 16)


# -------------------- S17: 4-7 India Part Spec --------------------
def slide_partspec_india():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-7.  India 부품 사양  |  QXI · BI3")

    data = [
        ["Project", "F/R", "Rod OD", "Inner Tube", "Outer Tube", "Valve", "Part No."],
        ["QXI",        "Front",  "Ø20",   "Ø30×1.03T",  "Ø48.6×2.5T", "MTV", "MET1L1801"],
        ["QXI",        "Rear",   "Ø12.4", "Ø25.4×1.18T","Ø38.5×1.27T","MTV", "MET1L1801"],
        ["BI3 (Type1)","Front",  "Ø18",   "Ø30×1.03T",  "Ø45×2.5T",   "MTV", "MEM501701/801"],
        ["BI3 (Type2)","Front",  "Ø20",   "Ø30×1.03T",  "Ø45×2.5T",   "MTV", "MEM501901/502001"],
        ["BI3",        "Rear",   "Ø12.4", "Ø25.4×1.18T","Ø38.5×1.27T","MTV", "MEM503101"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(3.3),
               data, col_widths=[1.7, 0.9, 1.5, 2.0, 2.2, 1.2, 2.8],
               font_size=11, header_size=11,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # feature box
    add_rect(s, Inches(0.5), Inches(5.3), Inches(6.0), Inches(1.7),
             fill=BG_LIGHT)
    add_text(s, Inches(0.7), Inches(5.4), Inches(5.8), Inches(0.4),
             "🇮🇳  인도 사양 특징",
             size=13, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(5.8), Inches(5.8), Inches(1.1),
             "• Rod 직경 Ø18 · Ø20 두 타입 운영 (BI3 Dual-spec)\n" +
             "• Outer Tube Ø45~48.6 · 중국 대비 경량화 설계\n" +
             "• 전 사양 MTV 밸브 (보급형 튜닝)",
             size=11, color=TEXT_DARK)

    add_rect(s, Inches(6.7), Inches(5.3), Inches(6.1), Inches(1.7),
             fill=TEAL)
    add_text(s, Inches(6.9), Inches(5.4), Inches(5.9), Inches(0.4),
             "💡  차별화 포인트",
             size=13, bold=True, color=LIME)
    add_text(s, Inches(6.9), Inches(5.8), Inches(5.9), Inches(1.1),
             "• Venue (QXI): 소형 SUV 전용 최경량 사양\n" +
             "• Elite (BI3): 해치백 Dual-spec 대응력 입증\n" +
             "• 향후 Creta India 확대 시 플랫폼 공용 가능",
             size=11, color=WHITE)
    add_footer(s, 17)


# -------------------- S18: 4-8 Brazil Part Spec --------------------
def slide_partspec_brazil():
    s = add_slide()
    add_top_bar(s, 4, "테네코에서의 현대 공급 프로그램",
                "4-8.  Brazil 부품 사양  |  BR2 · SU2")

    data = [
        ["Project", "F/R", "Rod OD", "Inner Tube", "Outer Tube", "Valve", "Part No."],
        ["BR2 (HB20)",   "Front", "Ø20",    "Ø30×1.27T", "Ø45×2.5T",  "MTV", "MEM702301"],
        ["BR2 (HB20)",   "Rear",  "Ø12.4",  "Ø25.4×1.2T","Ø38.1×1.2T","RV+", "MEM702501"],
        ["SU2 (Creta)",  "Front", "Ø20",    "Ø30×1.27T", "Ø48.6×2.5T","MTV", "ME3579301/F401/F501/F601"],
        ["SU2 (Creta)",  "Rear",  "Ø12.4",  "Ø30×1.27T", "Ø45.3×1.2T","RV+", "MEM966301/966401"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(2.8),
               data, col_widths=[1.7, 0.9, 1.4, 2.0, 2.2, 1.2, 2.9],
               font_size=11, header_size=11,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # key finding
    add_rect(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
             fill=ACCENT)
    add_text(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.5),
             "🟠  브라질 전용: Rear에 RV+ 밸브 채택 → 남미 도로환경 맞춤 튜닝 (중국·인도와 차별화)",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 3 insight cards
    cards = [
        ("내경 튜브 두께 ↑", "Ø30×1.27T",
         "Inner tube 두께 1.03T → 1.27T 상향\n남미 노면 대응 내구성 강화"),
        ("RV+ 밸브 채용", "Rear 사양",
         "Ride Valve Plus\n중속 감쇠력 민감도 개선"),
        ("Part No. 다중화", "4개 변종",
         "SU2 Front ME3579301/F401\nF501/F601 동시 운영"),
    ]
    y = Inches(5.55)
    card_w = Inches(3.98)
    for i, (title, val, note) in enumerate(cards):
        x = Inches(0.5) + (card_w + Inches(0.12)) * i
        add_rect(s, x, y, card_w, Inches(1.45), fill=BG_LIGHT)
        add_rect(s, x, y, Inches(0.1), Inches(1.45), fill=TEAL)
        add_text(s, x + Inches(0.25), y + Inches(0.08), card_w - Inches(0.3),
                 Inches(0.3), title, size=12, bold=True, color=TEAL_DEEP)
        add_text(s, x + Inches(0.25), y + Inches(0.38), card_w - Inches(0.3),
                 Inches(0.4), val, size=16, bold=True, color=TEAL,
                 font=FONT_EN)
        add_text(s, x + Inches(0.25), y + Inches(0.82), card_w - Inches(0.3),
                 Inches(0.6), note, size=10, color=TEXT_MID)
    add_footer(s, 18)


# -------------------- S19: Section 5 divider --------------------
def slide_section5():
    s = add_slide()
    add_section_divider(s, 5, "경쟁사 비교 내용",
                        "Competitor Analysis  |  ZF · Mando · Bilstein vs. Tenneco", 19)


# -------------------- S20: 5-1 Competitor Overview --------------------
def slide_competitor_overview():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-1.  경쟁사 개요  |  4대 주요 공급사 프로파일")

    comps = [
        ("MANDO", "만도",
         "Korea", "한국 현지 Tier-1",
         "현대/기아 메인 공급사. 가격·대응력 우위. 2-valve 등 신기술 대응 어려움.",
         TEAL, WHITE),
        ("ZF",  "ZF Friedrichshafen",
         "Germany", "글로벌 Tier-1",
         "Ioniq6·Genesis Entry ECS 공급. 고성능 인지도 최상. 가격 高, Air Shock 미지원.",
         TEAL_DARK, WHITE),
        ("BILSTEIN", "Bilstein (thyssenkrupp)",
         "Germany", "프리미엄 전문",
         "Genesis High-end (Air Shock+ECS) 공급. 2-valve MacPherson 부재.",
         TEAL_DEEP, WHITE),
        ("S&T", "S&T Motiv",
         "Korea", "한국 2nd Tier-1",
         "국내 쇼크업소버 2번째 공급사. 가격 우위. 고성능·신기술 실적 제한적.",
         ACCENT, WHITE),
        ("TENNECO", "Tenneco (Monroe)",
         "USA", "CES Business Hyundai",
         "2-valve MacPherson·Double Wishbone 대응. 가격 합리. 고성능 실적 확보 필요.",
         LIME, TEAL_DEEP),
    ]
    # 3 columns layout (2 rows: 3 + 2)
    y0 = Inches(1.85)
    card_h = Inches(2.5)
    card_w = Inches(4.05)
    for i, (en, full, country, tier, desc, bg, textc) in enumerate(comps):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + (card_w + Inches(0.12)) * col
        y = y0 + (card_h + Inches(0.12)) * row
        # card
        add_rect(s, x, y, card_w, card_h, fill=BG_LIGHT)
        # colored header bar
        add_rect(s, x, y, card_w, Inches(0.7), fill=bg)
        add_text(s, x + Inches(0.25), y + Inches(0.05), Inches(3), Inches(0.35),
                 en, size=18, bold=True, color=textc, font=FONT_EN)
        add_text(s, x + Inches(0.25), y + Inches(0.38), Inches(3), Inches(0.3),
                 full, size=10, color=textc, italic=True)
        # country badge
        add_text(s, x + card_w - Inches(1.7), y + Inches(0.2),
                 Inches(1.5), Inches(0.3), country,
                 size=11, bold=True, color=textc, align=PP_ALIGN.RIGHT,
                 font=FONT_EN)
        # tier
        add_text(s, x + Inches(0.25), y + Inches(0.85),
                 card_w - Inches(0.4), Inches(0.35),
                 f"포지션  |  {tier}", size=11, bold=True, color=TEAL_DEEP)
        # desc
        add_text(s, x + Inches(0.25), y + Inches(1.25),
                 card_w - Inches(0.4), card_h - Inches(1.4),
                 desc, size=11, color=TEXT_DARK)

    add_footer(s, 20)


# -------------------- S21: 5-2 Supplier Matrix --------------------
def slide_competitor_matrix():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-2.  공급사 비교 매트릭스  |  HMC R&D 관점")

    data = [
        ["공급사", "Office\n(한국)", "Plant\n(한국)", "R&D\n(한국)", "강점 (Strength)", "약점 (Weakness)"],
        ["Mando",   "○", "○", "○",
         "빠른 대응력 · 가격 경쟁력",
         "2-valve 등 신기술 부재 / 신규 차종 개발 대응 한계"],
        ["ZF",      "○", "○", "○",
         "고성능 라인 인지도 최상 · 개발 대응력 우수",
         "상대적으로 높은 가격 · Air Shock + ECS 조합 부재"],
        ["Bilstein","○", "✕", "○",
         "고성능 · 프리미엄 브랜드 이미지",
         "상대적으로 높은 가격 · 2-valve MacPherson 부재"],
        ["S&T",     "○", "○", "△",
         "국내 2nd 벤더 · 가격 우위",
         "고성능·신기술 개발 실적 제한적"],
        ["Tenneco", "○", "✕", "✕",
         "가격 합리성 (vs ZF/Bilstein) · 2025.4 유럽 Golf 시승 호평",
         "현대 고성능 프로젝트 개발 실적 없음"],
    ]
    tbl = make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.2),
                     data, col_widths=[1.3, 1.0, 1.0, 1.0, 4.0, 4.0],
                     font_size=10, header_size=11,
                     first_col_fill=TEAL_DEEP, first_col_color=WHITE)
    # highlight Tenneco row (last = row 5)
    try:
        last_row = tbl.rows[5]
        for c in last_row.cells:
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xDD)
    except Exception:
        pass

    # bottom takeaway
    add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.85),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
             "💡  Tenneco 포지셔닝",
             size=13, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
             "Plant/R&D 국내 부재 극복이 핵심 과제 → 공격 포인트: 기술 차별성(2-valve) + 가격 합리성 + 유럽 트랙레코드",
             size=11, color=TEAL_DEEP)
    add_footer(s, 21)


# -------------------- S22: 5-3 N Brand ECS Status --------------------
def slide_nbrand_ecs():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-3.  Hyundai N Brand ECS 공급 현황")

    data = [
        ["Eng", "Platform", "Model", "Damper Type", "Supplier", "핵심 사항"],
        ["ICE","B/C 세단","i20 N",           "Conventional",      "Mando", "N.A."],
        ["ICE","B/C 세단","i30 N / Veloster N","Semi-active",     "Mando", "Mando ECS 최대치 튜닝 → 추가 마진 부재"],
        ["ICE","B/C 세단","Avante N",         "Semi-active",      "Mando", "Mando ECS — dual-valve 미지원 한계 드러남"],
        ["EV", "EV GMP", "Ioniq 5",          "Semi-active",       "Mando", "Ioniq 6 개발 단계에서 한계 표면화"],
        ["EV", "EV GMP", "Ioniq 6",          "Semi-active",       "ZF",    "Mando 대체 채택 · HMC 유럽 공급사 편견 해소"],
    ]
    tbl = make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(3.6),
                     data, col_widths=[0.9, 1.6, 2.4, 1.8, 1.4, 4.2],
                     font_size=11, header_size=11,
                     first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)
    # highlight Ioniq6 row (last)
    try:
        r = tbl.rows[5]
        for c in r.cells:
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xDD)
    except Exception:
        pass

    # Takeaway
    add_rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.5),
             fill=BG_LIGHT)
    add_rect(s, Inches(0.5), Inches(5.55), Inches(0.12), Inches(1.5),
             fill=TEAL)
    add_text(s, Inches(0.8), Inches(5.65), Inches(12), Inches(0.4),
             "🎯  시사점  |  Mando 독점 체제의 균열",
             size=13, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.8), Inches(6.05), Inches(12), Inches(1.0),
             "• Ioniq 6 ZF 채택은 2-valve 기술 확보 + 유럽계 공급사 수용성 확대의 이중 시그널\n" +
             "• Tenneco CVSA2 (2-valve + MacPherson 대응)는 Mando가 막지 못한 영역을 정확히 공략\n" +
             "• N Brand에서 입증된 기술은 Genesis 브랜드로 수직 확산되는 HMC 개발 프로세스 존재",
             size=11, color=TEXT_DARK)
    add_footer(s, 22)


# -------------------- S23: 5-4 Genesis Brand ECS Status --------------------
def slide_genesis_ecs():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-4.  Hyundai Genesis Brand ECS 공급 현황")

    data = [
        ["Eng", "Platform", "Model", "Damper Type", "Supplier", "핵심 사항"],
        ["ICE","Sedan","G70",        "Conventional", "Mando",                        "N.A."],
        ["ICE","Sedan","G80 (ICE/EV)","Semi-active", "ZF (Entry) / Bilstein (High)", "Air Shock + ECS — 프리미엄 전략"],
        ["ICE","Sedan","G90",         "Semi-active", "ZF (Entry) / Bilstein (High)", "High-end Bilstein Air Shock+ECS"],
        ["ICE","SUV",  "GV70 (ICE/EV)","Semi-active","ZF (Entry) / Bilstein (High)", "Air Shock + ECS — SUV 프리미엄"],
        ["ICE","SUV",  "GV80",         "Semi-active","ZF (Entry) / Bilstein (High)", "플래그십 SUV"],
        ["EV", "EV GMP","GV60",        "Semi-active","Mando",                        "Ioniq 5 리소스 공용"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(3.9),
               data, col_widths=[0.9, 1.2, 2.0, 1.5, 2.4, 4.3],
               font_size=11, header_size=11,
               first_col_fill=BG_SOFT, first_col_color=TEAL_DEEP)

    # Strategy note
    add_rect(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.2),
             fill=TEAL)
    add_text(s, Inches(0.7), Inches(5.92), Inches(12), Inches(0.4),
             "🏆  HMC 프리미엄 전략",
             size=13, bold=True, color=LIME)
    add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.8),
             "BMW/Benz/Bentley 대응 목적의 Air Suspension 도입 → 공급가 아닌 품질 우선으로 ZF·Bilstein 선택.\n" +
             "Tenneco가 CVSA2 기술 + 프리미엄 사양 확보 시 G80/GV70 Entry-Level ZF 대체 가능성 존재.",
             size=11, color=WHITE)
    add_footer(s, 23)


# -------------------- S24: 5-5 N.America Competitor Plants --------------------
def slide_na_plants():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-5.  북미 경쟁사 Plant 분포  |  ZF · Mando · S&T")

    # === LEFT: North America map with markers ===
    map_x = Inches(0.4)
    map_y = Inches(1.75)
    map_w = Inches(7.2)
    map_h = Inches(5.1)
    add_na_map(s, map_x, map_y, map_w, map_h)

    # Plants on NA map: (name, brand, rel_x, rel_y, color)
    # positions are rough normalized coords within the map box
    na_plants = [
        ("ZF — Detroit, MI",       "ZF",    0.56, 0.42, TEAL_DEEP),
        ("ZF — Gainesville, GA",   "ZF",    0.60, 0.58, TEAL_DEEP),
        ("Mando — Opelika, AL",    "Mando", 0.55, 0.60, TEAL),
        ("Mando — Monterrey, MX",  "Mando", 0.40, 0.82, TEAL),
        ("ZF — Arteaga, MX",       "ZF",    0.38, 0.80, TEAL_DEEP),
        ("S&T — Michigan, USA",    "S&T",   0.52, 0.44, ACCENT),
        ("Tenneco — Michigan, USA","Tenneco",0.54, 0.46, LIME),
    ]
    # draw pins
    for i, (name, brand, rx, ry, color) in enumerate(na_plants):
        cx = map_x + Emu(int(map_w * rx))
        cy = map_y + Emu(int(map_h * ry))
        # pin circle
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - Inches(0.14), cy - Inches(0.14),
                                 Inches(0.28), Inches(0.28))
        pin.shadow.inherit = False
        pin.fill.solid()
        pin.fill.fore_color.rgb = color
        pin.line.color.rgb = WHITE
        pin.line.width = Pt(1.5)
        add_text(s, cx - Inches(0.14), cy - Inches(0.14),
                 Inches(0.28), Inches(0.28),
                 str(i+1), size=8, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)

    # === RIGHT: Legend list ===
    add_rect(s, Inches(7.9), Inches(1.75), Inches(5.1), Inches(0.45),
             fill=TEAL_DEEP)
    add_text(s, Inches(7.9), Inches(1.75), Inches(5.1), Inches(0.45),
             "주요 북미 거점 (#번호 = 지도 마커)",
             size=12, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # legend rows
    legend_plants = [
        ("1", "ZF",     "Detroit, MI (USA)",   "Powertrain / Chassis HQ", TEAL_DEEP),
        ("2", "ZF",     "Gainesville, GA",     "Chassis / Suspension",    TEAL_DEEP),
        ("3", "Mando",  "Opelika, AL",         "Shock Absorber Plant",    TEAL),
        ("4", "Mando",  "Monterrey, MX (Arteaga)","Shock Absorber",       TEAL),
        ("5", "ZF",     "Arteaga, MX",         "Cluster Plant",           TEAL_DEEP),
        ("6", "S&T",    "Michigan, USA",       "Chassis Component",       ACCENT),
        ("7", "Tenneco","Michigan, USA",       "R&D + HQ (참고)",         LIME),
    ]
    y = Inches(2.3)
    for i, (num, brand, loc, role, color) in enumerate(legend_plants):
        yy = y + Inches(0.57*i)
        # number pin
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(7.95), yy + Inches(0.09),
                                 Inches(0.32), Inches(0.32))
        pin.shadow.inherit = False
        pin.fill.solid()
        pin.fill.fore_color.rgb = color
        pin.line.fill.background()
        add_text(s, Inches(7.95), yy + Inches(0.09), Inches(0.32),
                 Inches(0.32),
                 num, size=10, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        # brand
        add_rect(s, Inches(8.38), yy + Inches(0.09), Inches(0.95),
                 Inches(0.32), fill=color)
        add_text(s, Inches(8.38), yy + Inches(0.09), Inches(0.95),
                 Inches(0.32),
                 brand, size=9, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        # loc
        add_text(s, Inches(9.4), yy, Inches(3.5), Inches(0.25),
                 loc, size=10, bold=True, color=TEAL_DEEP)
        add_text(s, Inches(9.4), yy + Inches(0.25), Inches(3.5),
                 Inches(0.25),
                 role, size=9, color=TEXT_MID, italic=True)

    add_footer(s, 24)


# -------------------- S25: 5-6 S.America Competitor Plants --------------------
def slide_sa_plants():
    s = add_slide()
    add_top_bar(s, 5, "경쟁사 비교 내용",
                "5-6.  남미 경쟁사 Plant 분포  |  ZF · Mando · S&T")

    # === LEFT: South America map with markers ===
    map_x = Inches(0.4)
    map_y = Inches(1.75)
    map_w = Inches(7.2)
    map_h = Inches(5.1)
    add_sa_map(s, map_x, map_y, map_w, map_h)

    sa_plants = [
        ("ZF — São Bernardo, BR",  "ZF",    0.68, 0.32, TEAL_DEEP),
        ("ZF — Sorocaba, BR",      "ZF",    0.62, 0.38, TEAL_DEEP),
        ("Mando — Limeira, BR",    "Mando", 0.66, 0.36, TEAL),
        ("ZF — Arg. San Francisco","ZF",    0.36, 0.75, TEAL_DEEP),
        ("S&T — Curitiba, BR",     "S&T",   0.58, 0.48, ACCENT),
        ("Tenneco — Cotia, BR",    "Tenneco",0.62, 0.30, LIME),
    ]
    for i, (name, brand, rx, ry, color) in enumerate(sa_plants):
        cx = map_x + Emu(int(map_w * rx))
        cy = map_y + Emu(int(map_h * ry))
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - Inches(0.14), cy - Inches(0.14),
                                 Inches(0.28), Inches(0.28))
        pin.shadow.inherit = False
        pin.fill.solid()
        pin.fill.fore_color.rgb = color
        pin.line.color.rgb = WHITE
        pin.line.width = Pt(1.5)
        add_text(s, cx - Inches(0.14), cy - Inches(0.14),
                 Inches(0.28), Inches(0.28),
                 str(i+1), size=8, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)

    # === RIGHT: Legend ===
    add_rect(s, Inches(7.9), Inches(1.75), Inches(5.1), Inches(0.45),
             fill=TEAL_DEEP)
    add_text(s, Inches(7.9), Inches(1.75), Inches(5.1), Inches(0.45),
             "주요 남미 거점 (#번호 = 지도 마커)",
             size=12, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    legend_plants = [
        ("1", "ZF",     "São Bernardo, Brazil", "HQ / Production",        TEAL_DEEP),
        ("2", "ZF",     "Sorocaba, Brazil",     "Chassis / 전장",          TEAL_DEEP),
        ("3", "Mando",  "Limeira, Brazil",      "Shock Absorber",          TEAL),
        ("4", "ZF",     "San Francisco, ARG",   "Chassis / Driveline",     TEAL_DEEP),
        ("5", "S&T",    "Curitiba, Brazil",     "Chassis Component",       ACCENT),
        ("6", "Tenneco","Cotia, Brazil",        "현지 거점 (참고)",        LIME),
    ]
    y = Inches(2.3)
    for i, (num, brand, loc, role, color) in enumerate(legend_plants):
        yy = y + Inches(0.57*i)
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(7.95), yy + Inches(0.09),
                                 Inches(0.32), Inches(0.32))
        pin.shadow.inherit = False
        pin.fill.solid()
        pin.fill.fore_color.rgb = color
        pin.line.fill.background()
        add_text(s, Inches(7.95), yy + Inches(0.09), Inches(0.32),
                 Inches(0.32),
                 num, size=10, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_rect(s, Inches(8.38), yy + Inches(0.09), Inches(0.95),
                 Inches(0.32), fill=color)
        add_text(s, Inches(8.38), yy + Inches(0.09), Inches(0.95),
                 Inches(0.32),
                 brand, size=9, bold=True,
                 color=(TEAL_DEEP if color == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_text(s, Inches(9.4), yy, Inches(3.5), Inches(0.25),
                 loc, size=10, bold=True, color=TEAL_DEEP)
        add_text(s, Inches(9.4), yy + Inches(0.25), Inches(3.5),
                 Inches(0.25),
                 role, size=9, color=TEXT_MID, italic=True)

    # bottom note
    add_rect(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.15),
             fill=LIME)
    add_footer(s, 25)


# -------------------- S26: Section 6 divider --------------------
def slide_section6():
    s = add_slide()
    add_section_divider(s, 6, "CVSA2 프로그램",
                        "CVSA2 Program  |  Next-Gen 2-Valve Damper for Hyundai", 26)


# -------------------- S27: 6-1 CVSA2 Background --------------------
def slide_cvsa2_background():
    s = add_slide()
    add_top_bar(s, 6, "CVSA2 프로그램",
                "6-1.  CVSA2 프로그램 배경")

    # Left: What
    add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.2),
             fill=BG_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.55),
             fill=TEAL)
    add_text(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(0.55),
             "CVSA2 란?", size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(2.5), Inches(5.6), Inches(3.4),
             "Continuously Variable Semi-Active 2nd Gen\n\n" +
             "• Tenneco 2세대 반능동형 감쇠 제어 쇼크업소버\n" +
             "• 2-Valve (Dual Valve) 시스템 채택\n" +
             "• MacPherson Strut · Double Wishbone 양방 대응\n" +
             "• 차량 주행 상태에 따른 감쇠력 실시간 가변\n" +
             "• ZF · Mando 1-valve 대비 응답성/승차감 우위",
             size=12, color=TEXT_DARK)

    # Right: Why Hyundai needs it
    add_rect(s, Inches(6.7), Inches(1.8), Inches(6.1), Inches(4.2),
             fill=TEAL_DEEP)
    add_text(s, Inches(6.9), Inches(1.95), Inches(5.8), Inches(0.4),
             "왜 현대에 필요한가?", size=14, bold=True, color=LIME)
    bullets = [
        "Mando 1-valve: 감쇠 세팅 한계 — Ioniq 6에서 ZF로 전환된 이유",
        "차세대 HMC N·Genesis 라인업은 2-valve 필수 요구",
        "ZF 2-valve: 가격 高 · Air Shock 조합 부재",
        "Bilstein 2-valve: MacPherson 대응 제품 없음",
        "→  Tenneco CVSA2가 4사 중 유일하게 모든 조건 충족",
    ]
    y = Inches(2.5)
    for i, b in enumerate(bullets):
        yy = y + Inches(0.62*i)
        add_rect(s, Inches(6.9), yy + Inches(0.08),
                 Inches(0.18), Inches(0.18), fill=LIME)
        add_text(s, Inches(7.2), yy, Inches(5.5), Inches(0.6),
                 b, size=11, color=WHITE, anchor=MSO_ANCHOR.TOP)

    # bottom summary
    add_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.3),
             "🎯  CVSA2 = Hyundai 차세대 Chassis의 '미싱 퍼즐'",
             size=13, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.3),
             "기술 우위 + 가격 합리성 + 플랫폼 범용성 — 3대 요건을 동시 만족하는 유일한 솔루션",
             size=11, color=TEAL_DEEP)
    add_footer(s, 27)


# -------------------- S28: 6-2 CVSA2 Technical Strength --------------------
def slide_cvsa2_technical():
    s = add_slide()
    add_top_bar(s, 6, "CVSA2 프로그램",
                "6-2.  CVSA2 기술적 차별화  |  2-Valve System")

    # three comparison columns
    cols = [
        ("Mando 1-Valve", "기존",
         ["Single valve 구조",
          "저속/고속 감쇠 트레이드오프",
          "Ioniq 6에서 한계 드러남",
          "Mando ECS 성능 상한",
          "추가 마진 확보 어려움"],
         GRAY_LIGHT, TEXT_DARK),
        ("ZF 2-Valve", "경쟁",
         ["Dual valve 채용",
          "감쇠 분리 제어 가능",
          "Genesis/Ioniq 6 채택",
          "높은 공급 단가",
          "Air Shock 조합 불가"],
         TEAL_DARK, WHITE),
        ("Tenneco CVSA2", "우리",
         ["Dual valve (2-Valve)",
          "MacPherson + Double Wishbone",
          "가격 합리 (vs ZF/Bilstein)",
          "유럽 시승 호평 (2025.4)",
          "N Brand → Genesis 확산 가능"],
         LIME, TEAL_DEEP),
    ]

    col_w = Inches(4.0)
    col_h = Inches(5.0)
    y = Inches(1.8)
    for i, (name, tag, items, bg, textc) in enumerate(cols):
        x = Inches(0.5) + (col_w + Inches(0.2)) * i
        # card
        add_rect(s, x, y, col_w, col_h, fill=BG_LIGHT)
        # header
        add_rect(s, x, y, col_w, Inches(1.1), fill=bg)
        add_text(s, x, y + Inches(0.1), col_w, Inches(0.3),
                 f"[ {tag} ]", size=10, bold=True, color=textc,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(s, x, y + Inches(0.38), col_w, Inches(0.6),
                 name, size=18, bold=True, color=textc,
                 align=PP_ALIGN.CENTER, font=FONT_EN,
                 anchor=MSO_ANCHOR.MIDDLE)
        # items
        for j, it in enumerate(items):
            yy = y + Inches(1.35) + Inches(0.65*j)
            # check icon
            mark = "✗" if i == 0 else ("△" if i == 1 else "✓")
            mcol = RED if i == 0 else (ACCENT if i == 1 else GREEN)
            add_text(s, x + Inches(0.2), yy, Inches(0.35), Inches(0.4),
                     mark, size=16, bold=True, color=mcol,
                     align=PP_ALIGN.CENTER)
            add_text(s, x + Inches(0.55), yy + Inches(0.02),
                     col_w - Inches(0.6), Inches(0.6),
                     it, size=11, color=TEXT_DARK)

    # bottom tag
    add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.08),
             fill=LIME)
    add_footer(s, 28)


# -------------------- S29: 6-3 Demo Car Status --------------------
def slide_demo_car_status():
    s = add_slide()
    add_top_bar(s, 6, "CVSA2 프로그램",
                "6-3.  Demo Car 검증 현황  |  Target: Ioniq 6 N")

    # Target box
    add_rect(s, Inches(0.5), Inches(1.8), Inches(5.9), Inches(2.6),
             fill=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.4),
             "🎯  Target Vehicle", size=12, bold=True, color=LIME)
    add_text(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(0.7),
             "Ioniq 6 N",
             size=32, bold=True, color=WHITE, font=FONT_EN)
    add_text(s, Inches(0.7), Inches(3.2), Inches(5.5), Inches(0.4),
             "EV GMP Platform · High-Performance EV Sedan",
             size=11, color=LIME, italic=True)
    add_text(s, Inches(0.7), Inches(3.65), Inches(5.5), Inches(0.7),
             "HMC가 Ioniq 5에서 드러난 Mando ECS 한계를\n" +
             "가장 명확히 확인할 수 있는 고성능 N 브랜드 차량",
             size=11, color=WHITE)

    # Benchmark box (right)
    add_rect(s, Inches(6.6), Inches(1.8), Inches(6.2), Inches(2.6),
             fill=BG_LIGHT)
    add_text(s, Inches(6.8), Inches(1.9), Inches(6), Inches(0.4),
             "🧪  Benchmark 구성",
             size=12, bold=True, color=TEAL_DEEP)
    benches = [
        ("Base",  "ZF, 1-valve",      "현재 Ioniq 6 양산 사양",   TEAL_DARK),
        ("Test 1","Mando, 1-valve",   "현재 Ioniq 5 양산 사양",   ACCENT),
        ("Test 2","Tenneco, CVSA2",   "Tenneco 제안 신규 사양",    LIME),
    ]
    yy = Inches(2.4)
    for i, (label, spec, note, c) in enumerate(benches):
        y2 = yy + Inches(0.6*i)
        add_rect(s, Inches(6.8), y2, Inches(1.2), Inches(0.45), fill=c)
        add_text(s, Inches(6.8), y2, Inches(1.2), Inches(0.45),
                 label, size=11, bold=True,
                 color=(TEAL_DEEP if c == LIME else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_text(s, Inches(8.1), y2, Inches(2.5), Inches(0.22),
                 spec, size=11, bold=True, color=TEAL_DEEP, font=FONT_EN)
        add_text(s, Inches(8.1), y2 + Inches(0.22), Inches(4.5), Inches(0.22),
                 note, size=10, color=TEXT_MID)

    # Tenneco Strength in HMC view
    add_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.4),
             fill=BG_SOFT)
    add_rect(s, Inches(0.5), Inches(4.6), Inches(0.12), Inches(2.4),
             fill=LIME)
    add_text(s, Inches(0.8), Inches(4.7), Inches(12), Inches(0.4),
             "💪  HMC가 본 Tenneco 강점",
             size=13, bold=True, color=TEAL_DEEP)
    pts = [
        "CVSA2 2-valve는 MacPherson 및 Double Wishbone 양쪽 적용 가능 (특히 MacPherson에 강점)",
        "경쟁사 대비 합리적 공급가격 (HMC 내부 평가)",
        "유럽 현지 Golf 시승 이후 HMC 엔지니어 긍정 피드백 확보",
        "현 Mando 1-valve의 성능 한계를 HMC R&D도 인지",
        "Demo Car 테스트는 RFQ 검증 과정 — ZF와 동등 성능 확인 시 우위 확보 가능",
    ]
    for i, p in enumerate(pts):
        yy = Inches(5.15) + Inches(0.35*i)
        add_rect(s, Inches(0.9), yy + Inches(0.08),
                 Inches(0.18), Inches(0.18), fill=LIME)
        add_text(s, Inches(1.2), yy, Inches(11.5), Inches(0.4),
                 p, size=10, color=TEXT_DARK)
    add_footer(s, 29)


# -------------------- S30: 6-4 Demo Car Schedule --------------------
def slide_demo_car_schedule():
    s = add_slide()
    add_top_bar(s, 6, "CVSA2 프로그램",
                "6-4.  Demo Car 검증 일정  |  Namyang R&D")

    # timeline steps
    steps = [
        ("CW42", "Oct. 13~19", "샘플 선적",
         "Sample Shipment to HMC Namyang"),
        ("CW46", "Nov. 11~16", "전장 설치·커미셔닝",
         "Electric System Installation & Commissioning"),
        ("CW47~49","Nov. 17~Dec. 7","Demo-Car HW/SW 튜닝",
         "Performance Tuning & Verification"),
        ("End of\nJune 2026","By End of June 2026","개선안 수립",
         "Develop Improvement Plan → RFQ 준비"),
    ]
    box_w = Inches(2.95)
    box_h = Inches(3.5)
    y = Inches(1.9)
    for i, (cw, date, title, desc) in enumerate(steps):
        x = Inches(0.5) + (box_w + Inches(0.15)) * i
        # timeline card
        add_rect(s, x, y, box_w, box_h, fill=BG_LIGHT)
        # week tag
        add_rect(s, x, y, box_w, Inches(1.0), fill=TEAL_DEEP)
        add_text(s, x, y + Inches(0.1), box_w, Inches(0.55),
                 cw, size=20, bold=True, color=LIME,
                 align=PP_ALIGN.CENTER, font=FONT_EN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.65), box_w, Inches(0.3),
                 date, size=10, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # step num
        num = add_rect(s, x + box_w - Inches(0.6), y + Inches(1.1),
                       Inches(0.5), Inches(0.5), fill=LIME)
        add_text(s, x + box_w - Inches(0.6), y + Inches(1.1),
                 Inches(0.5), Inches(0.5),
                 f"{i+1:02d}", size=14, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        # content
        add_text(s, x + Inches(0.2), y + Inches(1.4),
                 box_w - Inches(0.4), Inches(0.8),
                 title, size=13, bold=True, color=TEAL_DEEP)
        add_text(s, x + Inches(0.2), y + Inches(2.3),
                 box_w - Inches(0.4), Inches(1.0),
                 desc, size=10, color=TEXT_MID, italic=True)

    # Arrow-like bottom bar
    add_rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.1),
             fill=LIME)

    # Bottom: location
    add_rect(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.15),
             fill=BG_SOFT)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
             "📍  검증 장소 및 체계",
             size=13, bold=True, color=TEAL_DEEP)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12), Inches(0.4),
             "• 샘플 선적지: 남양 HMC R&D (Korea)",
             size=11, color=TEXT_DARK)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "• 전담 조직: High Performance Vehicle Engineering Design Team · Chassis Parts Purchasing Team 1",
             size=11, color=TEXT_DARK)
    add_footer(s, 30)


# -------------------- S31: 6-5 CVSA2 Expansion Plan --------------------
def slide_cvsa2_expansion():
    s = add_slide()
    add_top_bar(s, 6, "CVSA2 프로그램",
                "6-5.  현대에서의 CVSA2 확대 방안")

    # 3-stage roadmap
    stages = [
        ("STAGE 1", "진입 / Foothold",
         "Ioniq 6 N (Demo 차량)",
         ["Demo Car 검증 및 RFQ 수주",
          "ZF 1-valve 성능 동등 + 2-valve 우위 입증",
          "N Brand 내 최초 Tenneco 수주 확보"],
         TEAL),
        ("STAGE 2", "확장 / Expansion",
         "N Brand 추가 모델 + Genesis Entry",
         ["Ioniq 5 N 후속 · Avante N FL",
          "G80 · GV70 Entry Level ECS 공급",
          "N Brand → Genesis 수직 확산"],
         TEAL_DARK),
        ("STAGE 3", "주류 / Mainstream",
         "Genesis High-end + HMC 일반 라인업",
         ["G90 · GV80 High-end Air Shock + CVSA2",
          "HMC 대량 모델 (Ioniq 시리즈) 점진 채택",
          "Tenneco Korea HMC 주력 공급사 포지셔닝"],
         TEAL_DEEP),
    ]
    y = Inches(1.85)
    stage_w = Inches(4.0)
    stage_h = Inches(4.65)
    for i, (num, kr, target, items, bg) in enumerate(stages):
        x = Inches(0.5) + (stage_w + Inches(0.15)) * i
        add_rect(s, x, y, stage_w, stage_h, fill=BG_LIGHT)
        # header
        add_rect(s, x, y, stage_w, Inches(1.1), fill=bg)
        add_text(s, x + Inches(0.25), y + Inches(0.12),
                 stage_w - Inches(0.3), Inches(0.35),
                 num, size=11, bold=True, color=LIME, font=FONT_EN)
        add_text(s, x + Inches(0.25), y + Inches(0.45),
                 stage_w - Inches(0.3), Inches(0.55),
                 kr, size=18, bold=True, color=WHITE)
        # target row
        add_rect(s, x, y + Inches(1.1), stage_w, Inches(0.6), fill=BG_SOFT)
        add_text(s, x + Inches(0.25), y + Inches(1.15),
                 stage_w - Inches(0.3), Inches(0.22),
                 "Target", size=9, bold=True, color=TEAL_DEEP,
                 font=FONT_EN)
        add_text(s, x + Inches(0.25), y + Inches(1.37),
                 stage_w - Inches(0.3), Inches(0.3),
                 target, size=11, bold=True, color=TEAL_DEEP)
        # items
        for j, it in enumerate(items):
            yy = y + Inches(1.85) + Inches(0.75*j)
            add_rect(s, x + Inches(0.25), yy + Inches(0.1),
                     Inches(0.15), Inches(0.15), fill=bg)
            add_text(s, x + Inches(0.5), yy,
                     stage_w - Inches(0.7), Inches(0.8),
                     it, size=10, color=TEXT_DARK)

    # bottom arrow
    add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.3),
             "2026 (진입) → 2027~2028 (확장) → 2029~ (주류화)",
             size=11, bold=True, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 31)


# -------------------- S32: Section 7 divider --------------------
def slide_section7():
    s = add_slide()
    add_section_divider(s, 7, "현대 공략 방법",
                        "Hyundai Attack Strategy  |  Technology × Relationship × Execution", 32)


# -------------------- S33: 7-1 Launching Process --------------------
def slide_launching_process():
    s = add_slide()
    add_top_bar(s, 7, "현대 공략 방법",
                "7-1.  Hyundai ECS Launching Process  |  3-Stage Gate")

    # 3 stages
    stages = [
        ("01", "Technology\nShowcase",
         "유럽 (독일·벨기에)\nTenneco ↔ Hyundai",
         "Vehicle driving\nTechnical exchange",
         "차량 시승 및 기술 교류\n고성능 브랜드 선호 VP 대응"),
        ("02", "Engineering\n& Development",
         "HKMC R&D (Korea, 화성)",
         "Demo car build\nPerformance test\nInternal validation",
         "High Performance Vehicle\nEngineering Design Team"),
        ("03", "Purchasing",
         "HKMC HQ (Korea, 서울)",
         "RFQ / RFI\nQuotation",
         "Chassis Parts\nPurchasing Team 1"),
    ]
    box_w = Inches(3.9)
    box_h = Inches(4.6)
    y = Inches(1.85)
    for i, (num, title, loc, activities, team) in enumerate(stages):
        x = Inches(0.5) + (box_w + Inches(0.25)) * i
        # body
        add_rect(s, x, y, box_w, box_h, fill=BG_LIGHT)
        # big number on top
        add_rect(s, x, y, box_w, Inches(1.1), fill=TEAL_DEEP)
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(1),
                 Inches(0.8),
                 num, size=44, bold=True, color=LIME, font=FONT_EN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.3), y + Inches(0.2),
                 box_w - Inches(1.5), Inches(0.8),
                 title, size=15, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        # location
        add_rect(s, x, y + Inches(1.1), box_w, Inches(0.6), fill=LIME)
        add_text(s, x + Inches(0.2), y + Inches(1.15),
                 box_w - Inches(0.3), Inches(0.55),
                 loc, size=10, bold=True, color=TEAL_DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        # activities
        add_text(s, x + Inches(0.3), y + Inches(1.9),
                 box_w - Inches(0.6), Inches(0.4),
                 "Activities", size=10, bold=True, color=TEAL_DEEP,
                 font=FONT_EN)
        add_text(s, x + Inches(0.3), y + Inches(2.25),
                 box_w - Inches(0.6), Inches(1.4),
                 activities, size=11, color=TEXT_DARK)
        # team
        add_line(s, x + Inches(0.3), y + Inches(3.65),
                 x + box_w - Inches(0.3), y + Inches(3.65),
                 color=GRAY_LINE)
        add_text(s, x + Inches(0.3), y + Inches(3.75),
                 box_w - Inches(0.6), Inches(0.3),
                 "Counterpart", size=10, bold=True, color=TEAL_DEEP,
                 font=FONT_EN)
        add_text(s, x + Inches(0.3), y + Inches(4.0),
                 box_w - Inches(0.6), Inches(0.6),
                 team, size=10, color=TEXT_MID, italic=True)

    # bottom flow note
    add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "Validation stage (N Performance) → Sourcing stage (Top-Down to Genesis) — N Brand 검증 제품은 Genesis로 수직 확산",
             size=11, bold=True, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 33)


# -------------------- S34: 7-2 공략 Framework --------------------
def slide_attack_framework():
    s = add_slide()
    add_top_bar(s, 7, "현대 공략 방법",
                "7-2.  공략 Framework  |  3-Axis Strategy")

    axes = [
        ("기술 (Technology)",
         "2-Valve CVSA2 기술 우위 확정",
         ["MacPherson·Double Wishbone 양 플랫폼 대응",
          "경쟁사 대비 dual-valve 유일성 입증",
          "Air Shock + ECS 통합 솔루션 로드맵"],
         TEAL),
        ("관계 (Relationship)",
         "HMC 핵심 인물 공략",
         ["VP Manfred Harrer : 유럽계 선호 적극 활용",
          "Porsche 레퍼런스 강조 (HMC 긍정 인식)",
          "남양 R&D 엔지니어 유럽 초청 · 시승 프로그램 확대"],
         TEAL_DARK),
        ("실행 (Execution)",
         "N Brand → Genesis 확산 경로",
         ["Ioniq 6 N RFQ 수주 (Stage 1)",
          "2026 Demo Car 개선안 6월 확정 → 2027 양산",
          "Tenneco Korea Plant/R&D 단계별 구축"],
         TEAL_DEEP),
    ]
    y = Inches(1.9)
    card_w = Inches(4.0)
    card_h = Inches(4.8)
    for i, (title, tagline, items, bg) in enumerate(axes):
        x = Inches(0.5) + (card_w + Inches(0.15)) * i
        # body
        add_rect(s, x, y, card_w, card_h, fill=BG_LIGHT)
        # header
        add_rect(s, x, y, card_w, Inches(1.2), fill=bg)
        add_text(s, x + Inches(0.3), y + Inches(0.15),
                 card_w - Inches(0.4), Inches(0.45),
                 f"AXIS {i+1:02d}", size=10, bold=True,
                 color=LIME, font=FONT_EN)
        add_text(s, x + Inches(0.3), y + Inches(0.45),
                 card_w - Inches(0.4), Inches(0.7),
                 title, size=18, bold=True, color=WHITE)
        # tagline
        add_rect(s, x, y + Inches(1.2), card_w, Inches(0.65), fill=BG_SOFT)
        add_text(s, x + Inches(0.3), y + Inches(1.22),
                 card_w - Inches(0.4), Inches(0.63),
                 tagline, size=12, bold=True, color=TEAL_DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        # items
        for j, it in enumerate(items):
            yy = y + Inches(2.05) + Inches(0.85*j)
            add_rect(s, x + Inches(0.3), yy + Inches(0.08),
                     Inches(0.18), Inches(0.18), fill=bg)
            add_text(s, x + Inches(0.6), yy,
                     card_w - Inches(0.8), Inches(0.9),
                     it, size=11, color=TEXT_DARK)
    add_footer(s, 34)


# -------------------- S35: 7-3 Key Action Items --------------------
def slide_key_actions():
    s = add_slide()
    add_top_bar(s, 7, "현대 공략 방법",
                "7-3.  Key Action Items  |  즉시 실행 과제")

    # Table of actions
    data = [
        ["#", "과제", "세부 내용", "Owner", "Timeline"],
        ["1", "Ioniq 6 N Demo Car RFQ 수주",
         "CW47~49 HW/SW 튜닝 → 2026.6 개선안 → RFQ 제출",
         "Tenneco Korea + CES Engineering", "2026 Q1~Q2"],
        ["2", "유럽 Technology Showcase 확대",
         "남양 엔지니어/HMC 매니저 유럽 초청 시승 (Golf/Porsche 포함)",
         "Tenneco EU + Sales", "2026 Q2~Q3"],
        ["3", "VP Manfred Harrer 전담 커뮤니케이션",
         "Porsche 레퍼런스 중심 Executive Meeting 분기 1회",
         "Tenneco NA HQ Exec", "상시"],
        ["4", "N Brand → Genesis 확산 로드맵 제안",
         "Genesis G80/GV70 Entry ECS 2027 전환 시나리오 제시",
         "Tenneco CES Product Planning", "2026 Q3"],
        ["5", "Tenneco Korea 체제 구축",
         "Plant·R&D 부재 극복 — 단계별 현지화 계획 수립",
         "Tenneco Korea 경영진", "2026~2028"],
    ]
    make_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.5),
               data, col_widths=[0.5, 2.7, 4.6, 2.5, 2.0],
               font_size=10, header_size=11,
               first_col_fill=TEAL_DEEP, first_col_color=WHITE)

    # Success metric
    add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
             fill=LIME)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
             "🏁  Success Metric  |  2026년 말 Ioniq 6 N CVSA2 수주 확정 → 2027년 SOP → 2028 Genesis Entry 확대",
             size=12, bold=True, color=TEAL_DEEP, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 35)


# -------------------- S36: Section 8 divider --------------------
def slide_section8():
    s = add_slide()
    add_section_divider(s, 8, "Tenneco Korea의 단계별 운영 방안",
                        "Tenneco Korea Operation Roadmap  |  Phase 1 · 2 · 3", 36)


# -------------------- S37: 8-1 Phased Operation Roadmap --------------------
def slide_phased_roadmap():
    s = add_slide()
    add_top_bar(s, 8, "Tenneco Korea 단계별 운영 방안",
                "8-1.  단계별 운영 로드맵  |  Korea Local Capability")

    phases = [
        ("PHASE 1", "기반 확보", "2026",
         "Sales · Engineering 인력 보강",
         ["CES 영업 전담팀 + Technical Sales 인력 확충",
          "HMC 남양 R&D 상주 엔지니어 배치",
          "유럽 시승 및 기술 교류 운영 채널 가동"]),
        ("PHASE 2", "실행 확장", "2027~2028",
         "Local Engineering Center 설립",
         ["Tenneco Korea Engineering Center 개소 (위치 협의)",
          "Ioniq 6 N SOP 지원 및 차종 확장 대응",
          "Genesis Entry Level 사전 대응 R&D 수행"]),
        ("PHASE 3", "현지화 완성", "2029~",
         "Local Plant + R&D 운영",
         ["현대/기아 신차 개발 초기부터 Tier-1 참여",
          "국내 공장·R&D 동시 운영으로 Mando 수준 대응력 확보",
          "연간 5개 이상 HMC 프로그램 동시 공급 체제"]),
    ]
    y = Inches(1.85)
    card_w = Inches(4.0)
    card_h = Inches(4.7)
    for i, (ph, kr, yr, tagline, items) in enumerate(phases):
        x = Inches(0.5) + (card_w + Inches(0.15)) * i
        bg = [TEAL, TEAL_DARK, TEAL_DEEP][i]
        # body
        add_rect(s, x, y, card_w, card_h, fill=BG_LIGHT)
        # header
        add_rect(s, x, y, card_w, Inches(1.4), fill=bg)
        add_text(s, x + Inches(0.25), y + Inches(0.1),
                 card_w - Inches(0.4), Inches(0.3),
                 ph, size=11, bold=True, color=LIME, font=FONT_EN)
        add_text(s, x + Inches(0.25), y + Inches(0.4),
                 card_w - Inches(0.4), Inches(0.55),
                 kr, size=20, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), y + Inches(0.95),
                 card_w - Inches(0.4), Inches(0.4),
                 yr, size=12, bold=True, color=LIME, font=FONT_EN)
        # tagline
        add_rect(s, x, y + Inches(1.4), card_w, Inches(0.55), fill=BG_SOFT)
        add_text(s, x + Inches(0.25), y + Inches(1.4),
                 card_w - Inches(0.4), Inches(0.55),
                 tagline, size=12, bold=True, color=TEAL_DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        # items
        for j, it in enumerate(items):
            yy = y + Inches(2.15) + Inches(0.8*j)
            add_rect(s, x + Inches(0.3), yy + Inches(0.08),
                     Inches(0.15), Inches(0.15), fill=bg)
            add_text(s, x + Inches(0.55), yy,
                     card_w - Inches(0.7), Inches(0.85),
                     it, size=10, color=TEXT_DARK)
    # timeline arrow at bottom
    add_rect(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.25),
             fill=LIME)
    add_footer(s, 37)


# -------------------- S38: 8-2 Organization & Resource --------------------
def slide_org_resource():
    s = add_slide()
    add_top_bar(s, 8, "Tenneco Korea 단계별 운영 방안",
                "8-2.  조직 · 리소스 · 투자 계획")

    # Left: Organization
    add_rect(s, Inches(0.5), Inches(1.8), Inches(6.2), Inches(5.2),
             fill=BG_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.8), Inches(6.2), Inches(0.5),
             fill=TEAL)
    add_text(s, Inches(0.7), Inches(1.8), Inches(6), Inches(0.5),
             "🏢  조직 구성 (Target 2027)",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # org nodes
    org_items = [
        ("Korea Country Head", "총괄 · HMC Key Account",     "1명"),
        ("CES Sales & Marketing","영업·MKT·Technical Sales",   "6명"),
        ("Engineering Center", "남양 상주 + 자체 Lab",         "12명"),
        ("Program Mgmt (PMO)", "5개 HMC 프로젝트 동시 관리", "4명"),
        ("Quality & SQA",      "HMC 품질 기준 대응",           "3명"),
        ("Local Plant (Phase3)","브라질·중국 플랜트 Mirror",   "TBD"),
    ]
    y = Inches(2.4)
    for i, (role, detail, num) in enumerate(org_items):
        yy = y + Inches(0.68*i)
        add_rect(s, Inches(0.7), yy, Inches(0.1), Inches(0.55), fill=TEAL)
        add_text(s, Inches(0.95), yy, Inches(3.8), Inches(0.3),
                 role, size=12, bold=True, color=TEAL_DEEP)
        add_text(s, Inches(0.95), yy + Inches(0.28), Inches(3.8), Inches(0.3),
                 detail, size=10, color=TEXT_MID)
        add_rect(s, Inches(5.3), yy + Inches(0.08), Inches(1.2),
                 Inches(0.4), fill=LIME)
        add_text(s, Inches(5.3), yy + Inches(0.08), Inches(1.2),
                 Inches(0.4),
                 num, size=11, bold=True, color=TEAL_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)

    # Right: Investment & KPI
    add_rect(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(2.5),
             fill=TEAL_DEEP)
    add_text(s, Inches(7.1), Inches(1.9), Inches(5.7), Inches(0.4),
             "💰  투자 우선순위",
             size=13, bold=True, color=LIME)
    invest = [
        "[1] HMC 남양 상주 Engineering — 2026 Q1",
        "[2] Korea Technical Office (강남/판교) — 2026 Q2",
        "[3] CVSA2 시험/분석 장비 — 2026 Q3",
        "[4] Local Engineering Center — 2027 H1",
        "[5] Local Plant Feasibility — 2028",
    ]
    for i, v in enumerate(invest):
        yy = Inches(2.4) + Inches(0.34*i)
        add_rect(s, Inches(7.1), yy + Inches(0.08), Inches(0.15),
                 Inches(0.15), fill=LIME)
        add_text(s, Inches(7.35), yy, Inches(5.5), Inches(0.3),
                 v, size=10, color=WHITE)

    # KPI box
    add_rect(s, Inches(6.9), Inches(4.45), Inches(5.9), Inches(2.55),
             fill=LIME)
    add_text(s, Inches(7.1), Inches(4.55), Inches(5.7), Inches(0.4),
             "📊  핵심 KPI",
             size=13, bold=True, color=TEAL_DEEP)
    kpis = [
        ("2026",  "Ioniq 6 N RFQ 수주 · Demo 검증 완료"),
        ("2027",  "Genesis Entry ECS 1~2개 모델 수주"),
        ("2028",  "HMC 연매출 USD 100M 돌파"),
        ("2029",  "HMC Top-3 Chassis 공급사 진입"),
    ]
    for i, (yr, kp) in enumerate(kpis):
        yy = Inches(5.0) + Inches(0.48*i)
        add_rect(s, Inches(7.1), yy, Inches(0.9), Inches(0.4), fill=TEAL_DEEP)
        add_text(s, Inches(7.1), yy, Inches(0.9), Inches(0.4),
                 yr, size=11, bold=True, color=LIME,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_EN)
        add_text(s, Inches(8.1), yy + Inches(0.05), Inches(4.6),
                 Inches(0.35),
                 kp, size=10, color=TEAL_DEEP, bold=True)
    add_footer(s, 38)


# -------------------- S39: Section 9 divider --------------------
def slide_section9():
    s = add_slide()
    add_section_divider(s, 9, "Q & A",
                        "Questions & Discussion", 39)


# -------------------- S40: Closing / Thank You --------------------
def slide_closing():
    s = add_slide()
    # BG
    add_rect(s, 0, 0, SW, SH, fill=TEAL_DEEP)
    # lime accent
    add_rect(s, 0, Inches(2.8), SW, Inches(0.08), fill=LIME)

    add_text(s, 0, Inches(1.7), SW, Inches(1.2),
             "THANK YOU", size=80, bold=True, color=LIME,
             align=PP_ALIGN.CENTER, font=FONT_EN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, Inches(3.1), SW, Inches(0.6),
             "감사합니다", size=30, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, Inches(3.9), SW, Inches(0.5),
             "Tenneco × Hyundai Motor Group",
             size=16, color=LIME, italic=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, 0, Inches(4.5), SW, Inches(0.4),
             "Next-Gen Chassis Partnership · 2026 →",
             size=12, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_EN)

    # bottom ribbon
    add_rect(s, 0, Inches(6.6), SW, Inches(0.3), fill=LIME)
    add_text(s, 0, Inches(6.6), SW, Inches(0.3),
             "TENNECO — RIDE ENGINEERED FOR HYUNDAI",
             size=11, bold=True, color=TEAL_DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_EN)
    add_text(s, Inches(0.3), Inches(7.1), Inches(10), Inches(0.3),
             "General Business – Tenneco Confidential",
             size=9, color=TEXT_LIGHT, italic=True)
    add_text(s, Inches(12.7), Inches(7.1), Inches(0.55), Inches(0.3),
             "40", size=9, color=TEXT_LIGHT,
             align=PP_ALIGN.RIGHT, font=FONT_EN)


# ===== Run =====
if __name__ == "__main__":
    slide_cover()                       # 1
    slide_toc()                         # 2
    slide_section1()                    # 3
    slide_global_market()               # 4
    slide_section2()                    # 5
    slide_america_detail()              # 6
    slide_america_vs_competitor()       # 7
    slide_section3()                    # 8
    slide_tenneco_share()               # 9
    slide_section4()                    # 10
    slide_supply_overview()             # 11
    slide_china_project()               # 12
    slide_india_brazil_project()        # 13
    slide_program_volume()              # 14
    slide_partspec_china1()             # 15
    slide_partspec_china2()             # 16
    slide_partspec_india()              # 17
    slide_partspec_brazil()             # 18
    slide_section5()                    # 19
    slide_competitor_overview()         # 20
    slide_competitor_matrix()           # 21
    slide_nbrand_ecs()                  # 22
    slide_genesis_ecs()                 # 23
    slide_na_plants()                   # 24
    slide_sa_plants()                   # 25
    slide_section6()                    # 26
    slide_cvsa2_background()            # 27
    slide_cvsa2_technical()             # 28
    slide_demo_car_status()             # 29
    slide_demo_car_schedule()           # 30
    slide_cvsa2_expansion()             # 31
    slide_section7()                    # 32
    slide_launching_process()           # 33
    slide_attack_framework()            # 34
    slide_key_actions()                 # 35
    slide_section8()                    # 36
    slide_phased_roadmap()              # 37
    slide_org_resource()                # 38
    slide_section9()                    # 39
    slide_closing()                     # 40

    import os
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "Tenneco_Hyundai_Presentation_KR.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")

